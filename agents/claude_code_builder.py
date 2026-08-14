"""
Claude Code slide builder — uses claude-agent-sdk to spawn an actual Claude
Code agent per module. The agent composes the .pptx using PptxGenJS (for new
builds) or python-pptx (for targeted edits of existing decks), produces
matplotlib charts + graphviz flowcharts as needed, and iterates with PNG
vision feedback until the deck looks right.

Requires:
  - `claude-agent-sdk` (pip)
  - `claude` CLI on PATH (`npm install -g @anthropic-ai/claude-code`)
  - Node.js + pptxgenjs (`cd production/tools && npm install`)
  - LibreOffice for slide rendering (`soffice` on PATH)
  - ANTHROPIC_API_KEY in env
"""
from __future__ import annotations
import base64
import json
import os
import shutil
import tempfile
import traceback
from pathlib import Path
from typing import Optional

from ..tools import render_tool, code_exec_tool
from ..tools.pptxgen_tool import PptxGenSession
from ..tools.slide_layouts import get_palette

# ── Load design directive (injected into the builder's system prompt) ──
_DESIGN_DIRECTIVE = ""
_directive_path = Path(__file__).parent.parent / "skills" / "slide_design_directive.md"
if _directive_path.exists():
    _DESIGN_DIRECTIVE = _directive_path.read_text()

MAX_TURNS = 90

# The builder runs an agentic PptxGenJS tool loop (generate code, render, self-correct). Pin it to
# Opus 5 — Anthropic's flagship agentic-coding model, same $5/$25 price as Opus 4.8 and same-tier-
# or-better, so it should drive this loop at least as reliably as 4.8 did. The CLI accepts the full
# ID (verified: claude-agent-sdk 0.2.127 / CLI 2.1.219 resolves claude-opus-5 → subtype=success).
#
# Why NOT a lower tier: Sonnet 5 FAILED here — on this loop it thrashed (no-op "flush" calls, sleeps,
# throwing errors to inspect the pptx object, and a DESTRUCTIVE pptx.slides.splice(0,...) reset) and
# wrote corrupt non-ZIP files ("Builder produced corrupt file"). Stay on an Opus-tier model. The
# builder-prompt guardrails + the corrupt-output retry in build() are the safety net if any model
# misbehaves — worth watching the first real build after this change lands.
#
# Cost is captured from the SDK's exact total_cost_usd (see pricing.py), so the dashboard stays
# accurate automatically.
BUILDER_MODEL = "claude-opus-5"


# QA-derived design rules shared by BOTH the fresh-build and targeted-edit prompts. Each block
# exists because Visual QA kept re-flagging the defect class; keeping them here (not in one
# prompt only) is what stops first builds from making mistakes that a paid QA rebuild pass then
# has to fix. Any future QA-derived rule goes HERE, never into just one of the two prompts.
# (Plain text, no braces — this passes through the templates' .format() call.)
ANTI_DEFECT_DESIGN_RULES = """
BOXES / OUTLINED CARDS — USE SPARINGLY (avoid the "everything in a rectangle" look):
- A visible outlined box is a deliberate EMPHASIS device, not a default wrapper. Do NOT draw a
  bordered rectangle around every bullet, heading, or text block. Plain text on the background,
  separated by whitespace, reads cleaner and more modern — reach for an outline only to set ONE
  element apart (e.g. a single hero callout), at most one or two boxed elements per slide.
- NEVER nest a box in a box (this is the most common artifact): if an element already sits on a
  filled/colored card, do NOT also put a bordered rectangle inside or around it. One frame max.
- Bullet lists, titles, and the source/citation line should NOT be individually boxed.
- Prefer a soft FILL (or nothing) over a visible BORDER; if content is already grouped by position,
  it does not also need an outline.

HORIZONTAL RULES / DIVIDERS — keep them clear of the text:
- If you draw a horizontal rule or divider under a title, give it clear vertical clearance BELOW the
  title's lowest line — it must never cut through, touch, or overlap the title text. Titles wrap to
  two lines often; position the rule below the WRAPPED height, not the single-line height.
- Leave a margin between the rule and whatever follows it too (a table header row, the top node of a
  diagram, the first bullet) so the rule doesn't collide with the next element.
- On the CONTENT master (which already has a left accent bar), a title rule should START at the right
  edge of that accent bar and align to it — never run the rule through, across, or past the vertical
  bar.

TITLES — clean text, no template residue:
- A title is plain prose: never start it with stray punctuation (", Linear models" / "; Recap") and
  never leave raw markdown markers (**, ##, backticks) in any visible slide text.

MATH ON SLIDES — Unicode inline, make_equation for display; NEVER raw TeX as text:
- TeX/LaTeX markup must NEVER appear as literal slide text: no $...$, no \\frac, no ^{{...}} or
  _{{...}} braces, anywhere a learner can see (titles, bullets, labels, chart legends).
- INLINE math inside a sentence or title uses real Unicode characters: Greek letters directly
  (λ, ω, ζ, π), operators (× ÷ ≤ ≥ ≈ ≠ → ∞ ∑ ∫ √ ∂), and Unicode superscripts for simple
  powers (x², e⁻ᵏᵗ, 10⁻³). If an inline expression is too complex for that, simplify the
  sentence or move the expression into a display equation.
- DISPLAY equations (anything with stacked fractions, integrals with limits, multi-term
  derivations) are IMAGES: call make_equation (e.g. latex='\\frac{{dy}}{{dt}} = -ky') with the
  slide's text color, inspect the returned PNG, and place it sized by the returned aspect
  ratio. In matplotlib chart labels/legends, wrap math in $...$ so mathtext typesets it —
  a legend showing literal e^{{-x}} braces is a defect.
"""


SYSTEM_PROMPT_PPTXGEN = """You are the Syndara slide architect. You build a PowerPoint
presentation for a training course module. You have access to Read, Write, Bash
plus a set of custom tools that let you run PptxGenJS JavaScript code, render
matplotlib charts, render graphviz flowcharts, and view slide PNGs so you can
see your own work.

YOUR JOB, in order:
1. Read the provided module outline JSON from {outline_path}. If the outline
   has an "approved_slide_plan" field, the creator has already reviewed and
   signed off on a detailed markdown content plan — FOLLOW IT. The plan's
   `markdown` string is the authoritative source for slide content, layouts,
   visuals, and speaker notes.

   The plan separates "**On-slide content**" from "**Speaker notes**" for
   every slide. This is a HARD contract:
   - "On-slide content" is what the learner SEES. Follow the plan's layout and
     text for each slide. Be concise — no more words than the point needs, and
     not so few the slide is cryptic on its own. A visual-forward slide is a
     VISUAL with short labels; a text-forward slide (bullets, columns, or a table)
     carries fuller lines so it actually explains the idea. Match the plan: don't
     force a diagram onto a slide it built as text-forward, and don't turn a
     visual slide into paragraphs. Almost never a bare, unstructured wall of text
     (references excepted). On text-forward slides, bold the 1–3 KEY words per
     line in the ACCENT color (rich-text runs) for scannability — never bold whole
     lines or every word. CREATOR'S REQUEST WINS: if the course description or
     instructions ask for mostly-text, mostly-visual, verbose, or concise slides,
     honor that over these defaults.
   - "Speaker notes" is what the narrator SAYS: the full in-depth prose
     teaching, examples, citations, nuance. Copy ALL of this into
     slide.addNotes() VERBATIM. Do NOT paste any of the speaker notes
     onto the slide surface. The narration is the teaching; the slide is
     the visual anchor.
   - EXCEPTION — references/citations slides: the FIRST references slide's
     notes are ONE short sign-off sentence like "These are the references used
     in this module — thank you for listening!". ADDITIONAL references slides
     get NO notes at all (call addNotes('') — no narration; they're just more
     references). Never read URLs, citations, or source names aloud.
   - SOURCE ON FACT SLIDES: if a slide's plan **Data source** field is NOT 'N/A',
     render that citation text VERBATIM as a small line (~8pt, subtext color). If
     'N/A', render nothing. Use the plan's exact wording; don't reword. It's a
     caption/attribution, not slide body text.
     WHERE to place it:
       • If the cited data is shown in a CHART, GRAPH, or TABLE on the slide, put
         the citation DIRECTLY BENEATH that chart/table (like a caption) — NOT in a
         slide corner. (Web images likewise carry their caption directly beneath
         the image.)
       • Otherwise — a fact/figure stated in plain TEXT with no chart/table — place
         it in the free CORNER, picked from which overlays are active this build
         (each reserved-zone block appears in this prompt only when active):
           - presenter-video tile owns bottom-LEFT (~27% W × ~28% H); watermark
             owns bottom-RIGHT (x > 10", y > 6.5").
           - BOTH the "PRESENTER VIDEO — RESERVED ZONE" and "WATERMARK RESERVED
             ZONES" blocks present → TOP-RIGHT (clear of the title).
           - ONLY the watermark block → BOTTOM-LEFT.
           - ONLY the presenter-video block, OR NEITHER → BOTTOM-RIGHT.
     Never let the citation enter an ACTIVE reserved zone; keep it clear of content.
   When in doubt, follow the plan's intent: keep visual slides visual, and let
   text-forward slides carry the words they need to make sense on their own.
   The speaker notes carry the depth; the slide carries the point.
2. Compose the slide deck by calling run_pptxgen_code with PptxGenJS
   JavaScript code. Each call adds slides to the in-memory presentation.
   The file is saved automatically when the build finishes. Do NOT call
   pptx.writeFile() — it is handled for you.
3. THE PLAN OWNS THE CONTENT; YOU OWN THE LAYOUT. What goes on each slide —
   its text, whether it has a visual, and exactly what that visual shows — is
   already decided in the plan. Render it faithfully; do NOT add, drop, swap, or
   upgrade content. Your judgment is for LAYOUT only: spacing, positioning,
   sizing, alignment, and the PptxGenJS/tool mechanics to realize what's specified.
   VISUALS ARE PLAN-DRIVEN. The plan's **Visual elements** section specifies
   exactly what visual each slide needs (type + detailed description). You
   MUST generate every visual the plan calls for — a slide with
   Visual type != 'none' that ships as text-only bullets is a FAILURE. But the
   converse is equally binding: a slide with **Type:** none is TEXT-ONLY — never
   add an icon, image, chart, or decorative shape the plan didn't ask for.
   Build the specified visual FIRST, then add the minimal on-slide text around it.
   Pick the right tool for the job:
     - **Linear process flows and step sequences**: build with MANUAL
       PptxGenJS shapes (addShape + addText). Native vector objects are
       crisper than any rendered PNG and match the slide palette exactly.
       See SERPENTINE LAYOUT below for chains longer than 5 steps. When you
       draw arrows/connectors between nodes, keep every arrow in one diagram
       the SAME weight and color, and place any edge label (e.g. Yes/No)
       clear of the line so the line never crosses its own label.
     - **Simple charts** (bar, line, pie, doughnut with clean data): use
       PptxGenJS NATIVE CHARTS via slide.addChart() in run_pptxgen_code.
       These are editable in PowerPoint and look crisp at any zoom.
     - **Complex charts** (heatmap, violin, multi-axis, heavy annotations):
       make_chart (matplotlib).
     - **Complex branching diagrams, architecture diagrams, decision trees
       with many interconnected nodes**: make_d2_diagram. Use `layout:
       "dagre"` for flow, `"elk"` for strict hierarchy / right-angle
       edges, `"tala"` for complex architecture. Reserve D2 for diagrams
       where automatic edge routing between many nodes adds real value —
       NOT for simple linear chains.
     - **Sequence diagrams, entity-relationship diagrams, gantt charts,
       state machines, class diagrams**: make_mermaid_diagram. This is
       Mermaid's sweet spot.
     - **Icons** (for icon+text rows, grid headers, stats): make_icon.
       Renders react-icons to PNG. Look up valid names with find_icon('concept')
       first — don't guess. Use accent-colored icons at 0.5-0.8 in.
       CONTRAST: the icon color must stand out against whatever is behind it —
       NEVER a near-background color (e.g. white/very-light on a light tile), or
       it vanishes. Default to the deck accent color on light fills.
       NEVER use emoji characters (🔒 🔍 etc.) in slide text — they are
       fixed-color, low-res, and look unprofessional. Use make_icon instead.
     - make_flowchart (graphviz) is legacy — only fall back to it if D2
       fails.
   DIAGRAM IMAGE ASPECT RATIO — NEVER DISTORT:
   D2/Mermaid/matplotlib return the rendered PNG with exact pixel dimensions
   and aspect ratio. When inserting the image on a slide, you MUST preserve
   the original aspect ratio. Compute one dimension from the other:
     w = targetH * aspect;  // or  h = targetW / aspect;
   NEVER set both w and h independently — that stretches or squishes the
   diagram, making text unreadable. If the diagram doesn't fit the slide
   at a readable size while maintaining its aspect ratio, don't use
   D2/Mermaid — build it manually with PptxGenJS shapes instead.
   BEFORE inserting any diagram image:
     (a) Look at the returned PNG — is it legible? labels clipped? colors
         wrong? empty nodes? If anything looks off, regenerate with
         different source or parameters.
     (b) Use the returned aspect ratio to compute w and h as above.
   Insert the rendered PNG via run_pptxgen_code using slide.addImage().
   Keep surrounding on-slide text minimal — the visual is the point.
4. After every ~3 slides, call render_slide to view the PNG of a slide you
   built. If something looks broken (text clipped, chart scaled wrong,
   elements overlapping, wrong colors, too much text), fix it. Iterate
   until clean.
5. Set speaker notes on every content slide using the plan's "Speaker notes"
   content verbatim (or near-verbatim). This is non-negotiable.
6. When done, respond with "DONE: built N slides" (no more tool calls).

VARIABLES IN SCOPE (every run_pptxgen_code call):
  pptx      — PptxGenJS instance (add slides, charts, masters)
  style     — palette object with # prefixed hex colors + font names
  s         — SAME palette but colors ALREADY STRIPPED of # — USE THIS
              e.g. s.bg is 'F7F8FC', s.accent is '4361EE', s.title_font is 'Calibri'
  c(hex)    — helper that strips # from a string (for edge cases)
  PptxGenJS — the PptxGenJS class (for chart type constants)
  NOTE: these are the ONLY globals in run_pptxgen_code — there is NO require,
  fs, process, or import. Do not try to read the filesystem from this tool.

PRE-DEFINED SLIDE MASTERS (auto-created during init — use immediately):
  'CONTENT'    — palette bg + left accent bar motif
  'TITLE_DARK' — dark bg (accent color) + bottom accent strip
  'BLANK'      — palette bg, no decorations

GENERATED IMAGE FILES (charts / icons / diagrams):
  make_chart, make_icon, and make_d2_diagram each return "saved: <out_path>" —
  the exact path you passed in. You ALWAYS know your image paths (you chose
  them), so you never need to `ls` the images directory or probe the filesystem
  to find them. Insert directly with sl.addImage({{ path: '<out_path>' }}).

PptxGenJS QUICK REFERENCE:
  // Content slide using pre-defined master:
  let slide = pptx.addSlide({{ masterName: 'CONTENT' }});
  slide.addText('Title', {{ x: 1.0, y: 0.6, w: 11.33, h: 0.7, fontSize: 36,
      bold: true, color: s.text, fontFace: s.title_font, align: 'left' }});

  // Dark title slide:
  let title = pptx.addSlide({{ masterName: 'TITLE_DARK' }});
  title.addText('Course Title', {{ x: 1.0, y: 2.8, w: 11.33, h: 1.5,
      fontSize: 44, bold: true, color: s.bg, fontFace: s.title_font,
      align: 'center' }});

  // Rich text (mixed formatting):
  slide.addText([
      {{ text: 'Bold part', options: {{ bold: true, fontSize: 18 }} }},
      {{ text: ' normal part', options: {{ fontSize: 18 }} }},
    ], {{ x: 1.0, y: 1.6, w: 10, h: 3, color: s.text, fontFace: s.body_font }});

  // Shapes:
  slide.addShape('rect', {{ x: 0.8, y: 0.7, w: 0.08, h: 0.5,
      fill: {{ color: s.accent }}, line: {{ width: 0 }} }});
  slide.addShape('roundRect', {{ x: 1, y: 6.6, w: 2, h: 0.35,
      fill: {{ color: s.surface }}, line: {{ color: s.subtext, width: 1 }},
      rectRadius: 0.1 }});
  // For circles / number badges use 'ellipse' (a square w==h gives a circle).
  // CRITICAL: the valid preset names are 'ellipse', 'rect', 'roundRect' — NOT
  // 'oval', 'circle', or 'rectangle'. PptxGenJS passes an unknown name straight
  // through as prst="...", which is an invalid OOXML geometry: PowerPoint shows
  // a "repair" dialog and the shape renders as a broken diagonal line.
  slide.addShape('ellipse', {{ x: 1.5, y: 2, w: 0.5, h: 0.5,
      fill: {{ color: s.accent }}, line: {{ width: 0 }} }});
  // CONTRAST (text on filled shapes / table cells): the text color MUST contrast
  // with the fill it sits on. NEVER use s.accent or s.accent2 for text on a shape
  // or cell filled with that same color — it renders invisible (e.g. orange text
  // on an orange header cell). Rule: on an accent or dark fill use color: s.bg
  // (light text); on a light fill (s.bg / s.surface / s.highlight) use color:
  // s.text (dark text).

  // Text inside a shape (diagram label) — ALWAYS use fit:'none' + align:'center':
  // ALL text inside shapes/boxes/cards MUST be center-aligned to prevent
  // mixed alignment when text wraps to a second line.
  slide.addText('Step 1', {{ x: 1, y: 2, w: 2, h: 0.8,
      shape: pptx.shapes.ROUNDED_RECTANGLE, rectRadius: 0.15,
      fill: {{ color: s.accent }}, line: {{ width: 0 }},
      fontSize: 14, color: s.bg, fontFace: s.body_font,
      align: 'center', valign: 'middle', fit: 'none' }});

  // Images:
  slide.addImage({{ path: '/absolute/path/to/image.png', x: 2, y: 1.5, w: 9, h: 5 }});
  // Speaker notes:
  slide.addNotes('Full speaker notes here — verbatim from plan');

  // Native chart (editable in PowerPoint — prefer for simple data):
  slide.addChart(pptx.charts.BAR, [
      {{ name: 'Series', labels: ['A','B','C'], values: [12,20,18] }},
    ], {{ x: 1.5, y: 1.8, w: 10, h: 4.5,
      showLegend: true, legendPos: 'b',
      chartColors: [s.accent, s.accent2],
      catAxisLabelFontSize: 10, valAxisLabelFontSize: 10 }});
  // Chart types: pptx.charts.BAR, LINE, PIE, DOUGHNUT, SCATTER, AREA, RADAR

  CRITICAL RULES:
  - Use `s.accent` not '#4361EE' — the `s` object already has # stripped.
  - Do NOT reuse options objects across addText/addShape calls —
    PptxGenJS mutates them. Create fresh object literals each time.
  - Do NOT call pptx.writeFile() — the file is saved automatically.
  - All dimensions are in inches (native PptxGenJS unit).
  - If you get a word-count warning in the response, trim that slide's text.

═══ GOLDEN EXAMPLE — complete content slide ═══
  let sl = pptx.addSlide({{ masterName: 'CONTENT' }});
  // Title with accent bar already from master
  sl.addText('Key Findings', {{ x: 1.0, y: 0.6, w: 11.33, h: 0.7,
      fontSize: 32, bold: true, color: s.text, fontFace: s.title_font }});
  // 3-stat callout (hero numbers)
  const stats = [
      {{ num: '94%', label: 'Accuracy' }},
      {{ num: '3.2s', label: 'Avg Response' }},
      {{ num: '50K+', label: 'Users Served' }},
  ];
  stats.forEach((st, i) => {{
      const cx = 1.0 + i * 4.0;
      sl.addText(st.num, {{ x: cx, y: 2.2, w: 3.5, h: 1.5,
          fontSize: 64, bold: true, color: s.accent,
          fontFace: s.title_font, align: 'center' }});
      sl.addText(st.label, {{ x: cx, y: 3.8, w: 3.5, h: 0.8,
          fontSize: 14, color: s.subtext,
          fontFace: s.body_font, align: 'center' }});
  }});
  sl.addNotes('These findings represent our Q4 analysis across...');

═══ GOLDEN EXAMPLE — dark title slide ═══
  let ts = pptx.addSlide({{ masterName: 'TITLE_DARK' }});
  ts.addText('Module Title', {{ x: 1.0, y: 2.5, w: 11.33, h: 1.5,
      fontSize: 44, bold: true, color: s.bg,
      fontFace: s.title_font, align: 'center' }});
  ts.addText('Subtitle goes here', {{ x: 1.0, y: 4.2, w: 11.33, h: 0.7,
      fontSize: 20, color: s.surface || s.bg,
      fontFace: s.body_font, align: 'center' }});

{visual_motif}

DARK TITLE SLIDES:
Use masterName 'TITLE_DARK' for title and conclusion slides. The
background color is `s.title_bg` (falls back to `s.accent` if unset).
Check the actual darkness of the title background to choose text color:
if the title_bg is dark, use s.bg for text; if it's light, use s.text.

SHAPE-BASED DIAGRAMS (hub-and-spoke, radial, process flows built with addShape):
When you build diagrams manually with PptxGenJS shapes instead of D2/Mermaid:
- LINES MUST GO FROM HUB EDGE TO SATELLITE EDGE: the #1 bug is lines that
  start near the hub center and end near the hub edge — crossing over the hub
  without reaching the satellites. To fix this, compute line endpoints from
  shape edges, not centers:
    // Hub center and satellite center
    hubCX = hub.x + hub.w/2;  hubCY = hub.y + hub.h/2;
    satCX = sat.x + sat.w/2;  satCY = sat.y + sat.h/2;
    // Angle between them
    angle = Math.atan2(satCY - hubCY, satCX - hubCX);
    // For circles: line starts at hub EDGE, ends at satellite EDGE
    lineX  = hubCX + Math.cos(angle) * (hub.w/2);
    lineY  = hubCY + Math.sin(angle) * (hub.h/2);
    lineW  = (satCX - Math.cos(angle) * (sat.w/2)) - lineX;
    lineH  = (satCY - Math.sin(angle) * (sat.h/2)) - lineY;
  For rectangles, clamp to the CENTER OF THE CORRECT FACE (not corner):
    function rectEdge(cx, cy, w, h, angle) {{
      // Returns the point on the rectangle edge at the given angle from center
      let dx = Math.cos(angle), dy = Math.sin(angle);
      let scaleX = (w/2) / Math.abs(dx || 0.001);
      let scaleY = (h/2) / Math.abs(dy || 0.001);
      let scale = Math.min(scaleX, scaleY);
      return {{ x: cx + dx * scale, y: cy + dy * scale }};
    }}
    let startPt = rectEdge(hubCX, hubCY, hub.w, hub.h, angle);
    let endPt   = rectEdge(satCX, satCY, sat.w, sat.h, angle + Math.PI);
  This ensures arrows connect at the center of each face, not corners
  or along the edge. Never eyeball line coordinates.

DIAGONAL LINE DIRECTION — CRITICAL:
  PptxGenJS draws lines from (x,y) to (x+w, y+h). Both w and h MUST be
  positive. To change direction, use flipH and/or flipV:
    - Top-left → bottom-right: flipH:false, flipV:false (default)
    - Top-right → bottom-left: flipH:true,  flipV:false
    - Bottom-left → top-right: flipH:false, flipV:true
    - Bottom-right → top-left: flipH:true,  flipV:true
  The arrowhead (line.endArrowType:'triangle') always points at the END
  of the line AFTER flipping. So to get an arrow from A to B:
    1. Set x = min(Ax, Bx), y = min(Ay, By)
    2. Set w = abs(Bx - Ax), h = abs(By - Ay)
    3. Set flipH = (Bx < Ax), flipV = (By < Ay)
    4. Set line.endArrowType:'triangle' — arrow points at B
  Example (arrow pointing up-right, from bottom-left to top-right):
    sl.addShape('line', {{ x: 2, y: 1, w: 3, h: 2,
        flipV: true,
        line: {{ color: s.text, width: 2, endArrowType: 'triangle' }} }});
- SPREAD SATELLITES OUT: place satellites far from the hub so lines have
  visible length. If the hub is in the center of the slide, satellites
  should be near the edges. Don't cluster them close to the hub.
- LAYOUT BEFORE LINES: place all shapes first, then calculate connecting
  lines from the final shape positions.
- LABELS INSIDE SHAPES: font size should be readable (12-16pt minimum).
  If the text doesn't fit, make the shape bigger — don't shrink the text.
  ALWAYS use fit:'none' on text inside diagram shapes to prevent
  PowerPoint from auto-shrinking text (which squishes it horizontally).
  Example:
    sl.addText('Label', {{ x: 2, y: 3, w: 2, h: 1, fontSize: 14,
        fit: 'none', align: 'center', valign: 'middle',
        color: s.text, fontFace: s.body_font }});
  Size the shape (w/h) to comfortably contain the text at the chosen
  font size. For circles and squares in a sequence, keep each shape
  large enough for its label — don't cram long text into tiny shapes.
- WHEN TO USE D2/MERMAID vs MANUAL SHAPES:
  Use D2/Mermaid for complex branching diagrams with many interconnected
  nodes where automatic edge routing adds value. Use manual PptxGenJS
  shapes for linear chains, simple process flows, and any diagram where
  you need precise control over layout or the slide aspect ratio.
  Manual shapes are native vectors (crisper, editable in PowerPoint,
  match the palette exactly) — prefer them when the layout is simple
  enough to position by hand.

SERPENTINE LAYOUT (for linear chains with more than 5 steps):
When a process/flow has more than ~5 steps in a chain, wrap it into a
serpentine (boustrophedon) pattern so it fits the slide without shrinking
text. Split the steps roughly in half (lean toward more on top if odd).

Horizontal serpentine (for LR flows):
  Row 1 (L→R):  [Step 1] → [Step 2] → [Step 3]
                                              ↓
  Row 2 (R→L):  [Step 6] ← [Step 5] ← [Step 4]
  - Top row flows left-to-right, bottom row flows right-to-left
  - A single vertical arrow connects the last node of row 1 to the
    first node of row 2 (they share the same x-position on the right)

Vertical serpentine (for TB flows):
  Col 1 (T→B):  [Step 1]    Col 2 (B→T):  [Step 4]
                    ↓                          ↑
                [Step 2]                   [Step 5]
                    ↓                          ↑
                [Step 3]  →→→→→→→→→→→→→    [Step 6]
  - Left column flows top-to-bottom, right column flows bottom-to-top
  - A single horizontal arrow connects the last node of col 1 to the
    first node of col 2 (they share the same y-position at the bottom)

Implementation:
  1. Place all shapes first (two rows/columns with consistent spacing)
  2. Draw arrows between consecutive shapes in each row/column
  3. Draw the single turn arrow connecting the two rows/columns
  4. Use the flipH/flipV rules from DIAGONAL LINE DIRECTION above for
     arrows that go right-to-left or bottom-to-top

ICON USAGE:
Use make_icon to render professional vector icons (react-icons library).
Insert at 0.5-0.8 in for inline, 1.0-1.5 in for hero. Great for
icon+text rows, grid cell headers, stat callouts, and process steps.
GET THE NAME RIGHT, FAST: call find_icon with ONE plain concept noun —
'gear', 'lock', 'people', 'sun' — NOT a descriptive phrase ('gear flow loop
cycle' matches nothing and just burns turns). Pick a name from its results and
pass it, with its icon_pack, to make_icon. Do NOT guess icon_names from memory.
Keep it TIGHT: at most two find_icon tries per icon — if the second still finds
nothing that fits, render a simple geometric mark via make_icon (or drop the
icon) and move on; never keep re-searching or guessing names across packs. When a
slide's tiles each need a distinct icon, look each up once — never a repeated
icon, numbered circle, or single-LETTER initial standing in for a real icon.
NEVER paste a Unicode symbol or icon-font glyph as an icon in slide text —
not emoji (🔒 🔍 💡 ⚡), and not dingbat/symbol-font characters either. The
PPTX-to-PNG renderer does not have those glyph fonts, so they come out as
empty "tofu" boxes (blank rectangles) in the final slide. The ONLY way to
place an icon is make_icon → insert the returned PNG.
CONTRAST: give every icon a color that clearly stands out against whatever
sits behind it. Never render an icon in a near-background color (e.g. a
white or very-light icon on a light tile) — it becomes invisible. Use the
deck accent color on light fills, and a light color only on dark fills.
""" + ANTI_DEFECT_DESIGN_RULES + """
{progressive_builds_section}TEXT ROTATION:
PptxGenJS supports `rotate` (degrees, clockwise) on any element:
  sl.addText('Label', {{ x: 1, y: 2, w: 3, h: 0.6, rotate: 90,
      fontSize: 14, color: s.text, fontFace: s.body_font }});
Use rotation when it enhances the design — vertical labels along an axis,
text following a diagonal arrow, sideways section markers, angled callouts.
Default is horizontal (no rotate property). Don't rotate body text or
speaker notes — only short labels and decorative elements.
For python-pptx targeted edits, set `shape.rotation = degrees` (float,
clockwise).

QUESTION SLIDES (Layout: question_slide):
When the plan has a question_slide, produce TWO physical slides:
  1. QUESTION slide: shows the question text + answer options (A, B, C, D)
     with NO correct answer highlighted. All options look identical.
     Speaker notes MUST start with this exact metadata line:
       [QUESTION|options:Option A text;Option B text;Option C text;Option D text|correct:N]
     where N is 0-based index of the correct option. Follow it with the
     regular narration on a new line.
  2. ANSWER REVEAL slide: IDENTICAL layout but the correct answer option
     is highlighted (accent color background or bold border). Add a small
     "✓" or checkmark next to it.
     Speaker notes MUST start with exactly: [ANSWER_REVEAL]
     Follow it with the explanation of why the answer is correct.
This creates a click-to-reveal animation when presenting the PPTX.
Both slides should have the same background and positioning — only the
highlight on the correct answer differs.

NATIVE CHARTS vs MATPLOTLIB:
- Simple bar/line/pie/doughnut → use slide.addChart() (editable in PPT)
- Complex viz (heatmap, violin, multi-axis, annotations) → make_chart

HALF-BLEED IMAGES:
For high-impact visual slides, bleed an image to one edge (x=0, full height)
with text on the other half. Creates visual variety.

DIAGRAM GUIDANCE:
- Charts: label axes, include units, use accent color for emphasis, white
  gridlines on dark bg. Minimum 600x400 px.
- Linear flows / step sequences: build with manual PptxGenJS shapes (see
  SERPENTINE LAYOUT above). D2/Mermaid only for complex branching diagrams
  where automatic edge routing between many nodes adds real value.
- Don't produce a chart if the data doesn't warrant it. Placeholder bar
  charts with no insight are worse than a well-designed text slide.

CRITICAL VISUAL WORKFLOW — read this carefully:
For EACH slide in the plan:
  1. Parse the "**Visual elements**" -> "**Type:**" field
  2. If Type is NOT 'none', you MUST produce the visual. Either:
     - Build it with manual PptxGenJS shapes (for linear flows, step
       sequences, and simple layouts — this is preferred for crispness),
     - OR call a diagram tool (D2/Mermaid/matplotlib) for complex visuals
       that benefit from automatic layout, THEN insert the PNG.
     DO NOT skip this step. DO NOT write text-only slides when a visual
     is specified.
  3. The visual should be the dominant element on the slide (~9x5 inches).
  4. Text goes around/above the visual as small labels, not as the main
     content.
If you build a slide as text-only bullets when the plan says it needs a
visual, the entire build will be rejected. This wastes the creator's time
and yours. Generate the visual.

SPECIAL CASE — **Type:** animation: this slide's visual is a Manim ANIMATION rendered
separately and laid FULL-BLEED over the slide, covering it entirely. The player shows the
animation's own background before it plays, so this built slide is almost never seen — it
is ONLY a rare fallback (a failed render) and the static frame in the downloaded PPTX.
Build it BARE: use the TITLE_DARK master (its dark background best matches the animation)
and place ONLY the slide title on it. HARD RULES — no charts, no diagrams, no tables, no
bullets, no icons, no panels, no stats, no one-line takeaway; do NOT call make_chart /
make_d2_diagram / make_mermaid_diagram / any diagram tool for it, and do NOT approximate
the animation. Title on an empty dark background, full stop — anything more is wasted work
that just gets covered. It is EXEMPT from the "must have a visual" rule and must NOT be
flagged for having no diagram.

HARD LIMITS:
- Max {max_turns} total tool calls — a CEILING, not a target. Use only what the
  deck needs; a simple text deck should take just a few calls per slide. Do NOT
  pad with throwaway "test" slides, repeated re-renders of the same slide, or
  repeated icon lookups. Spend extra turns only where a real visual needs polish.
- Every slide with Visual type != 'none' MUST have a visual — either manual
  PptxGenJS shapes or a generated diagram/chart image. Text-only bullet
  slides when a visual was planned = failure.
- Keep on-slide text concise — as many words as the point needs, no filler. Visual
  slides get short labels; text-forward slides get fuller lines. Never a dense paragraph.
- Don't use Bash to install packages or download files. Everything you need
  is already available.
- Do not modify files outside the provided working directory.

DESIGN DIRECTIVE — follow these visual rules strictly:
{design_directive}
"""


# Teaching for executing the plan's reveal scripts. Injected into SYSTEM_PROMPT_PPTXGEN via the
# {progressive_builds_section} placeholder ONLY when the course has progressive builds enabled —
# a builds-off course gets a builder prompt in which the reveal-tag convention simply doesn't
# exist, so the model can't emit it (the plan upstream is equally build-free; and web_runner's
# explosion-time backstop strips anything that somehow slips through).
PROGRESSIVE_BUILDS_SECTION = """PROGRESSIVE BUILDS (executing the plan's reveal script — you do NOT choreograph):
A progressive build is a click-to-reveal on an ordinary slide (elements fade in as narrated) —
NOT a Manim Type: animation (that is the separate full-bleed video case below). Some plan slides
have Speaker notes written as a build SCRIPT: they start with [REVEAL] and contain inline markers
like [[1 | show: the three stage cards]]. The PLAN owns all build decisions. Your job is mechanical:
1. Author the slide ONCE at its complete, final layout — every element the script ever shows,
   all in their final positions (transients and both halves of a replacement included; on this
   authored slide they all render at once, and that is expected). Elements never move or restyle
   between beats; beats only show or hide elements (fade). Never reposition content for a
   partial state. When the script says a replacement sits "in place" of what it replaces,
   position it exactly there — the overlap is intentional staging.
2. Copy the script into slide.addNotes() VERBATIM — header, markers, descriptions, all of it.
3. Tag each animated element's objectName with its visibility window over the beat numbers:
     objectName:'reveal:2'       appears at beat 2, stays to the end
     objectName:'reveal:2-2'     visible ONLY during beat 2 (a transient)
     objectName:'reveal:1-1,3'   in at beat 1, out at 2, back from 3 on (re-entry)
   Read the windows off the script's descriptions: "[[2 | show: tip; hide: green note]]" means
   the tip's window starts at 2 and the green note's window ends at 1 (closed: '1-1' — unless a
   later marker shows it again, then '1-1,<that beat>'). Untagged elements are the base, visible
   throughout. Several elements sharing a window = they animate together.
4. Descriptions name elements loosely ("the ledger image") — map them to the shapes you created.
   If the plan's script is malformed (markers not 1,2,3,…, a description referencing nothing you
   can identify, an empty segment), build the slide STATIC with clean notes instead — a static
   slide is always acceptable; a broken build is not. Do not invent, add, drop, or reorder beats.
A slide whose plan notes have no [REVEAL] is a normal static slide — never add a build yourself.
(Downstream, each marker becomes one extra physical slide with the not-yet-shown elements hidden;
the markers never appear in the downloaded deck's notes, so copy them without worry.)

"""

SYSTEM_PROMPT_PPTX = """You are the Syndara slide architect. You edit an existing
PowerPoint presentation for a training course module. You have access to Read,
Write, Bash plus a set of custom tools that let you run python-pptx code,
render matplotlib charts, render graphviz flowcharts, and view slide PNGs so
you can see your own work.

YOUR JOB (TARGETED EDIT MODE):
You are editing specific slides in an existing deck using python-pptx. Use
run_pptx_code to open and modify the presentation. The `prs` (Presentation),
`pptx_path`, `Inches`, `Pt`, `Emu`, `RGBColor`, `PP_ALIGN`, `MSO_SHAPE` are
all in scope.

1. Call read_summary with slide_index=N to inspect the specific slide you
   need to edit. This returns full shape details (text, positions, sizes).
2. Use run_pptx_code to modify that slide. The tool automatically returns
   the structure of any modified slides — you do NOT need to call
   read_summary again after editing.
3. If a visual needs regenerating, use diagram tools then insert via
   run_pptx_code.
4. Call render_slide to verify each edited slide visually.
5. When done, respond with "DONE: edited N slide(s)".

EFFICIENCY — read this carefully:
- run_pptx_code returns the full structure of every slide you modified.
  Use print() in your code to output any values you need — they're captured
  and returned. Do NOT write workaround scripts to inspect slides.
- read_summary(slide_index=N) gives you detailed per-slide info. Use it
  BEFORE editing to understand shape layout. After editing, the return
  value of run_pptx_code already tells you the result.
- You should need at most 2-3 tool calls per slide: read_summary → edit →
  render_slide. If you're using more, you're doing it wrong.
- To REMOVE a shape it is ONE line: `sh._element.getparent().remove(sh._element)`.
  NEVER escalate to raw zip/lxml/XML surgery, external repair scripts, or
  reloading the whole deck to strip a class of element — that rabbit hole cost a
  prior run 25 minutes and several dollars. Connectors/lines are fine to keep;
  if one looks off, recolor/resize/reposition it, or delete just that one shape
  with the one-liner above.
- Fix ONLY the slides the feedback flagged. Do NOT sweep the whole deck hunting
  for elements to remove. If a flagged slide isn't fixed in ~5 calls, regenerate
  its visual as a single make_d2_diagram/make_chart image and move on.

Speaker notes: `notes_slide = slide.notes_slide; notes_slide.notes_text_frame.text = "..."`

TEXT ROTATION:
To rotate a shape or text box, set `shape.rotation = degrees` (float,
clockwise). Use for vertical labels, angled callouts, or text along
diagonal arrows. Default is 0 (horizontal). Only rotate short labels
and decorative elements — never body text.

QUESTION SLIDES:
If you're editing a question_slide pair (question + answer reveal), remember
these are TWO consecutive physical slides: the first shows the question with
unhighlighted options, the second shows the same layout with the correct
answer highlighted. Edit both slides to keep them in sync.
The question slide's speaker notes must start with:
  [QUESTION|options:A text;B text;C text;D text|correct:N]
The answer reveal slide's speaker notes must start with:
  [ANSWER_REVEAL]

LAYOUT GUIDANCE:
- 13.33 x 7.5 in (widescreen 16:9).
- Use the style palette's hex colors for bg/surface/accent/text/subtext.

TEXT IN SHAPES — CRITICAL:
- ALWAYS disable auto-shrink on text frames inside shapes. Without this,
  PowerPoint squishes text horizontally to fit, making it look low-res:
    from pptx.enum.text import MSO_AUTO_SIZE
    tf = shape.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True
- All text inside shapes/boxes/cards must be center-aligned:
    from pptx.enum.text import PP_ALIGN
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
- If text doesn't fit, make the shape bigger — never let it auto-shrink.
- Minimum font size: 10pt. Never go below this.
""" + ANTI_DEFECT_DESIGN_RULES + """
HARD LIMITS:
- Max {max_turns} total tool calls.
- Only modify slides you are told to modify. Follow structural rules in
  the user prompt (it will tell you whether add/delete is allowed).
- Do not modify files outside the provided working directory.

DESIGN DIRECTIVE — follow these visual rules strictly:
{design_directive}
"""


async def build_slides_with_claude_code(
    outline: dict,
    output_dir: str,
    style: str = "syndara",
    reviewer_feedback: Optional[dict] = None,
    max_turns: int = MAX_TURNS,
    watermark_info: Optional[dict] = None,
) -> str:
    """
    Produce a .pptx for the given outline via Claude Code.

    Returns the path to the generated file. Raises RuntimeError on failure.
    """
    try:
        from claude_agent_sdk import (
            query, ClaudeAgentOptions, ResultMessage, AssistantMessage,
            tool, create_sdk_mcp_server,
        )
    except ImportError as e:
        raise RuntimeError(
            f"claude-agent-sdk not installed: {e}. "
            "`pip install claude-agent-sdk` and ensure the `claude` CLI is on PATH."
        ) from e

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = str(output_dir / "slides.pptx")

    # Scale the tool budget with deck size: observed usage is ~2-2.5 turns per
    # slide (compose + chart/icon renders + self-check renders). A fixed budget
    # silently truncates large decks.
    target_slides = outline.get("slide_count") or len(outline.get("slides") or []) or 0
    if not target_slides:
        # In the plan-driven creator flow, slide_count/slides are empty — the authoritative count
        # is the number of "## Slide N" entries in the approved plan markdown. Without this the
        # budget never scales and 30-40 slide modules die on error_max_turns.
        _plan_obj = outline.get("approved_slide_plan") or {}
        _plan_md = _plan_obj.get("markdown", "") if isinstance(_plan_obj, dict) else (
            _plan_obj if isinstance(_plan_obj, str) else "")
        target_slides = _plan_md.count("\n## Slide ") + (1 if _plan_md.startswith("## Slide ") else 0)
    if target_slides:
        # ~3.5 tool calls/slide: icon- and diagram-heavy slides each spend several calls
        # (find_icon + make_icon + code + render), so 2.5/slide starved big decks — a
        # ~40-slide module hit the cap mid-build and was dropped entirely. Budget for the
        # deck to FINISH; the builder stops early when done, so extra headroom costs nothing
        # on normal modules and only spends more on the genuinely large ones.
        # Scale to the deck: don't floor tiny decks at the big default (a 3-slide
        # deck at 90 turns invites over-exploration — throwaway test slides, repeat
        # renders, icon-hunting). ~3.5 calls/slide + fixed setup, min 30.
        max_turns = max(30, int(target_slides * 3.5) + 20)
    outline_path = str(output_dir / "outline.json")
    images_dir = output_dir / "images"
    images_dir.mkdir(exist_ok=True)

    # ── Decide build vs targeted-edit mode ───────────────────────
    # If the Reviewer flagged specific slides (slide_index + issues +
    # suggestion) AND the .pptx from a prior pass still exists, don't wipe
    # it — only edit the flagged slides. Saves a full rebuild when most of
    # the deck is already clean.
    rf_slides = (reviewer_feedback or {}).get("slides") or []
    # Any per-slide status that asks for a change is a targeted edit. Critically this
    # includes "needs_diagram"/"needs_image": if only those are present, treating them
    # as un-actionable made is_targeted_edit False, which DELETED the deck and forced a
    # full rebuild over an "add a diagram to slide 7" note. They edit slides too.
    _ACTIONABLE_SLIDE_STATUSES = {"revise", "needs_diagram", "needs_image"}
    flagged = [
        s for s in rf_slides
        if isinstance(s, dict)
           and s.get("status") in _ACTIONABLE_SLIDE_STATUSES
           and "slide_index" in s
    ]
    pptx_exists = Path(pptx_path).exists()
    is_targeted_edit = bool(flagged) and pptx_exists

    # A targeted edit's budget scales with the FLAGGED slides, not the whole deck: a
    # one-slide fix on a 40-slide module was inheriting ~160 tool calls, which the prompt
    # itself warns "invites over-exploration". ~6 calls per flagged slide (read, edit,
    # chart, render, verify, margin) with a floor for setup/self-checks.
    if is_targeted_edit:
        # Retuned from max(12, n*6) after job 211: a QA fix session hit the 12-turn ceiling
        # mid-edit (read + chart regen + render + verify per slide) and errored the whole
        # pass — wasting the rebuild it was meant to be cheaper than. 24 + 8/slide still
        # sits far below the ~160-turn full-deck budget it exists to avoid.
        max_turns = min(max_turns, max(24, len(flagged) * 8))

    if not is_targeted_edit and pptx_exists:
        # Full rebuild (creator revision or no structured reviewer data) —
        # wipe the old deck so we don't append-duplicate slides.
        Path(pptx_path).unlink()

    # Serialize style palette for prompt
    palette = get_palette(style)
    if outline.get("plain_backgrounds"):
        palette["plain_backgrounds"] = True
    (output_dir / "style.json").write_text(json.dumps(palette, indent=2))
    Path(outline_path).write_text(json.dumps(outline, indent=2))

    # ── Decide engine: PptxGenJS for new builds, python-pptx for edits ──
    use_pptxgen = not is_targeted_edit

    # Start PptxGenJS session for new builds
    pptxgen_session = None
    if use_pptxgen:
        pptxgen_session = PptxGenSession(pptx_path, palette)

    if is_targeted_edit:
        allow_structural = (reviewer_feedback or {}).get("allow_structural_changes", False)
        flagged_summary_lines = []
        for s in flagged:
            idx = s["slide_index"]
            action = s.get("action", "revise")
            issues = ", ".join(s.get("issues", [])) or "unspecified"
            suggestion = s.get("suggestion", "").strip() or "(no suggestion)"
            action_label = f" [ACTION: {action.upper()}]" if action != "revise" else ""
            flagged_summary_lines.append(
                f"- prs.slides[{idx}] (user calls this \"slide {idx + 1}\"){action_label}: issues=[{issues}]\n    suggestion: {suggestion}"
            )
        global_fb = (reviewer_feedback or {}).get("global_feedback", "") or ""

        # When structural edits are allowed, keep the deck inside its length band.
        band_rule = ""
        if allow_structural and target_slides:
            from .slide_planner import slide_range_for
            _lo, _hi = slide_range_for(target_slides)
            _cur = None
            try:
                from pptx import Presentation as _Prs
                _cur = len(_Prs(pptx_path).slides)
            except Exception:
                pass
            band_rule = (
                f" DECK LENGTH BOUND: the deck must stay between {_lo} and {_hi} slides"
                + (f" (it currently has {_cur})" if _cur is not None else "")
                + f". You may add or delete as requested, but the final deck must "
                f"NOT have fewer than {_lo} or more than {_hi} slides — if a request "
                f"would cross that bound, make the edit in place instead of adding/removing."
            )
        structural_rule = (
            "STRUCTURAL CHANGES ALLOWED: You may delete or add slides as indicated "
            "by the [ACTION: DELETE] or [ACTION: ADD_AFTER] flags above. For DELETE: "
            "remove the slide entirely via `xml_slides = prs.slides._sldIdLst; "
            "xml_slides.remove(xml_slides[index])`. For ADD_AFTER: add a new slide "
            "after the indicated index. Adjust carefully — deletions shift indices."
            + band_rule
            if allow_structural else
            "Do NOT add or remove slides. Do NOT re-order slides. Only edit content in place."
        )

        prompt = f"""TARGETED EDIT — the .pptx from the previous pass already exists at {pptx_path}. Do NOT rebuild from scratch. Do NOT delete the file. Leave every slide NOT listed below completely untouched.

Style palette (same as before): {output_dir}/style.json
Module outline (for reference only): {outline_path}
Temporary image dir: {images_dir}

Slides needing revision (0-based indices matching prs.slides[N] — edit ONLY these):
{chr(10).join(flagged_summary_lines)}

Reviewer's overall note: {global_fb[:500]}

CRITICAL INDEXING RULE: The indices above are 0-based prs.slides[] indices.
Users count from 1, python-pptx counts from 0. The conversion is already done
for you above — "prs.slides[4]" IS what the user calls "slide 5". Use the
prs.slides[N] index EXACTLY as listed. Do NOT add or subtract. Do NOT touch
any slide not listed above.

{structural_rule}

Approach:
1. Call read_summary with slide_index=N to inspect the slide (includes speaker notes).
2. For each flagged slide, use run_pptx_code to make changes. `prs` is already
   loaded — do NOT reassign it with `prs = Presentation(...)`.
3. SPEAKER NOTES — handle ALL of these cases:
   a. Feedback is about speaker notes only → update notes, no visual changes needed.
   b. Feedback is about slide content only → change the slide, AND adapt the
      speaker notes if they reference the changed content (e.g., if you change
      "syndara.com" to "syndara.org" on the slide, fix it in the notes too).
   c. Feedback mentions both → update both.
   Access notes: `prs.slides[index].notes_slide.notes_text_frame.text`
   Notes are read aloud by TTS — they must match the on-slide content.
4. If a flagged slide needs a new/replacement visual, regenerate it with
   make_d2_diagram / make_mermaid_diagram / make_chart, inspect the
   returned PNG, then insert/replace it on that specific slide only.
5. After each slide edit, call render_slide on that index to verify the
   visual fix landed.
6. When all flagged slides are fixed, respond with "DONE: edited {len(flagged)} slide(s)".
"""
    else:
        feedback_block = ""
        if reviewer_feedback:
            _slides_fb = reviewer_feedback.get("slides") if isinstance(reviewer_feedback, dict) else None
            if _slides_fb:
                # Render as compact per-slide lines rather than dumping pretty-printed JSON and
                # slicing it at 3000 chars — that cut mid-structure and silently dropped the later
                # flagged slides. Cap on a whole-entry boundary with an explicit marker.
                _fb_lines = []
                for _s in _slides_fb:
                    _iss = ", ".join(str(i) for i in (_s.get("issues") or [])) or "(unspecified)"
                    _sug = (_s.get("suggestion") or "").strip()
                    _fb_lines.append(f"- slide {_s.get('slide_index')}: {_iss}" + (f" — {_sug}" if _sug else ""))
                    if len(_fb_lines) >= 60:
                        _fb_lines.append(f"...({len(_slides_fb) - 60} more flagged slides — see the outline)")
                        break
                _gfb = (reviewer_feedback.get("global_feedback") or "").strip()
                feedback_block = "\n\nPREVIOUS REVIEWER / CREATOR FEEDBACK — address these in this build:\n"
                if _gfb:
                    feedback_block += f"Overall: {_gfb}\n"
                feedback_block += "\n".join(_fb_lines)
            else:
                # Non-slide-structured feedback (free text): include it, capped on a line boundary.
                _raw = json.dumps(reviewer_feedback, indent=2)
                if len(_raw) > 3000:
                    _raw = _raw[:3000].rsplit("\n", 1)[0] + "\n  ...(truncated)"
                feedback_block = (
                    "\n\nPREVIOUS REVIEWER / CREATOR FEEDBACK — address these issues in this build:\n"
                    + _raw
                )

        max_words_per_slide = outline.get("max_words_per_slide") or 20

        # Extract approved slide plan markdown and embed directly in prompt
        # so the agent doesn't have to read a file and parse JSON first.
        plan_md = ""
        plan_obj = outline.get("approved_slide_plan") or {}
        if isinstance(plan_obj, dict):
            plan_md = plan_obj.get("markdown", "")
        elif isinstance(plan_obj, str):
            plan_md = plan_obj

        # Count slides from the plan markdown (## Slide N lines)
        slide_count = plan_md.count("\n## Slide ") + (1 if plan_md.startswith("## Slide ") else 0)
        if slide_count == 0:
            slide_count = len(outline.get("slides", [])) or 40

        # Presenter-video mode overlays a talking-head tile in the bottom-left corner of every
        # slide, so the builder must leave that region clear (see storage.video.PRESENTER_TILE).
        keepout_block = ""
        if outline.get("presenter_keepout"):
            keepout_block = (
                "\n═══════════════════════════════════════════════════════════\n"
                "PRESENTER VIDEO — RESERVED ZONE (hard layout constraint)\n"
                "═══════════════════════════════════════════════════════════\n"
                "A talking-head presenter video is overlaid in the BOTTOM-LEFT corner of EVERY slide.\n"
                "Keep the bottom-left ~28% of slide WIDTH and ~30% of slide HEIGHT completely clear —\n"
                "no text, bullets, charts, images, diagrams, or footers may sit there or they will be\n"
                "covered. Shift on-slide content up and/or to the right. This applies to every slide.\n"
            )

        prompt = f"""Build a {slide_count}-slide presentation for this course module.

Output: {pptx_path}
Style palette (read this file for exact hex colors): {output_dir}/style.json
Temporary image dir for charts/flowcharts: {images_dir}

Style name: "{style}"
Module: {outline.get('title', '')}

═══════════════════════════════════════════════════════════
APPROVED SLIDE PLAN — this is your authoritative blueprint.
Build EXACTLY what it specifies. Every slide, every visual.
═══════════════════════════════════════════════════════════

{plan_md or '(No detailed plan — use the outline at ' + outline_path + ')'}

═══════════════════════════════════════════════════════════
END OF PLAN
═══════════════════════════════════════════════════════════
{feedback_block}

INSTRUCTIONS — follow in order for EACH slide:
1. Read the plan section for this slide. Note the **Visual elements** type.
2. If the visual type is NOT 'none', generate the visual FIRST:
   - Linear flows / step sequences → manual PptxGenJS shapes (preferred —
     native vectors, crisper, palette-matched). Use serpentine wrap for >5
     steps.
   - Complex branching diagrams / architecture → make_d2_diagram, then
     insert PNG preserving aspect ratio.
   - Sequence / ER / gantt / state → make_mermaid_diagram, then insert
     PNG preserving aspect ratio.
   - Charts / stats → make_chart or native slide.addChart().
   For D2/Mermaid PNGs: inspect the returned image, regenerate if labels
   are clipped, and NEVER distort the aspect ratio when inserting.
3. Add the text items from "On-slide text" exactly as the plan lists them — concise, no filler.
4. Copy the full "Speaker notes" into slide.addNotes('...').
5. Call render_slide every 2-3 slides to verify.

A slide that the plan says needs a visual but you render as text-only
bullets is WRONG. You MUST generate the visual. This is non-negotiable.
{keepout_block}"""
        web_images = outline.get("web_images", {})
        # The key is dual-purpose: a BOOL enabled-flag on freshly-built outlines (read by the
        # exercise/worksheet image lookup) that the download step later replaces with the dict of
        # downloaded images. When a plan has no image slides the download step returns early and
        # the bool survives to here — iterating it crashed the whole module (204 m4).
        if not isinstance(web_images, dict):
            web_images = {}
        if web_images:
            img_lines = ["\nWEB IMAGES (pre-downloaded from the web — use these local paths):"]
            for heading, info in web_images.items():
                if not isinstance(info, dict):
                    continue
                if info.get("failed"):
                    img_lines.append(f"  {heading}: FAILED — {info.get('error', 'unknown')}. Build a diagram/icon set instead.")
                elif info.get("path"):
                    line = f"  {heading}: path={info['path']} ({info.get('width_px', 0)}x{info.get('height_px', 0)}px, aspect={info.get('aspect', 0)})"
                    if info.get("attribution"):
                        line += f" — attribution: {info['attribution']}"
                    if info.get("source"):
                        line += f" — source: {info['source']}"
                    img_lines.append(line)
                else:
                    img_lines.append(f"  {heading}: FAILED — no path available. Build a diagram/icon set instead.")
            prompt += "\n".join(img_lines)
            prompt += (
                "\n\nFor slides with a web image, use slide.addImage({ path: '<path>' }) "
                "with the local path above. Size to ~9x5 inches max, centered. "
                "REQUIRED — every web image MUST carry a formal source citation directly beneath it: "
                "a small caption (8pt italic, subtext color) reading 'Source: <attribution>, <source>' "
                "using the attribution and source URL given for that image above (drop a part only if "
                "it wasn't provided; NEVER invent a source). Do not place a web image without this "
                "citation line. If an image is marked FAILED, fall back to a generated "
                "visual (D2/Mermaid diagram, chart, or PptxGenJS shapes).\n"
                "\nIMAGE ASPECT RATIO — ALL IMAGES (web images, source images, photos, "
                "screenshots): NEVER stretch horizontally or vertically. Always scale "
                "uniformly — compute one dimension from the other to preserve the original "
                "aspect ratio. If the image doesn't fit at its natural ratio, make it smaller "
                "uniformly. If it would be too small, redesign the slide layout to accommodate "
                "it rather than distorting it.\n"
            )

    # ── Build MCP tools ──────────────────────────────────────────
    @tool(
        "run_pptx_code",
        "Run python-pptx code against the target .pptx file. prs (pre-opened Presentation), pptx_path, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_SHAPE are in scope. File is auto-saved after your code runs. Returns the structure of any modified slides so you can verify your changes without a separate read_summary call. print() output is also captured and returned.",
        {"code": str},
    )
    async def run_pptx_code_tool(args):
        result = code_exec_tool.run_pptx_code(args["code"], pptx_path)
        is_err = not result.get("success")
        if is_err:
            text = f"Error: {result.get('error')}"
        else:
            parts = [f"OK. Slide count: {result.get('slide_count')}"]
            if result.get("stdout"):
                parts.append(f"Output:\n{result['stdout']}")
            touched = result.get("slides_touched", [])
            if touched:
                parts.append(f"Modified slides: {json.dumps(touched, indent=1)}")
            text = "\n".join(parts)
        return {"content": [{"type": "text", "text": text}], "is_error": is_err}

    @tool(
        "run_pptxgen_code",
        "Run PptxGenJS JavaScript code to build slides. Variables in scope: "
        "`pptx` (PptxGenJS instance), `s` (palette with hex colors ALREADY "
        "stripped of #, e.g. s.bg='F7F8FC', s.accent='4361EE', "
        "s.title_font='Calibri'), `style` (same but with # prefix), "
        "`c(hex)` (strips #), `PptxGenJS`. Pre-defined masters: 'CONTENT' "
        "(bg + left accent bar), 'TITLE_DARK' (dark bg), 'BLANK'. "
        "Use pptx.addSlide({{ masterName: 'CONTENT' }}). Do NOT reuse "
        "options objects — PptxGenJS mutates them. Do not call writeFile(). "
        "STATE PERSISTS across calls: slides you added in earlier calls are still on the deck — the "
        "response tells you the current slide count. Build FORWARD, adding each slide once. NEVER "
        "reset or mutate the slide array to 'start over' (no pptx.slides.splice / pop / length=0 / "
        "reassignment) — that corrupts PptxGenJS's internal state and produces an unopenable file "
        "even though the save looks successful. If a slide came out wrong, fix THAT slide; do not "
        "reset the deck. Do NOT run no-op 'flush'/'verify'/'check' code, do NOT throw errors to "
        "inspect `pptx`, and do NOT sleep — none of that is needed and it just wastes turns.",
        {"code": str},
    )
    async def run_pptxgen_code_tool(args):
        result = pptxgen_session.run_code(args["code"])
        is_err = not result.get("success")
        if is_err:
            text = f"Error: {result.get('error')}"
        else:
            text = f"OK. Slide count now: {result.get('slide_count')}"
            warnings = result.get("warnings")
            if warnings:
                text += "\n⚠️ " + "\n⚠️ ".join(warnings)
        return {"content": [{"type": "text", "text": text}], "is_error": is_err}

    def _diagram_result(render_result: dict, kind: str) -> dict:
        """Shared return shape for all diagram-rendering tools. On success the
        rendered PNG is returned as a vision content block so the agent can
        actually see what it made before inserting into the deck. Dimensions
        are included so the agent can do aspect-ratio math."""
        if not render_result.get("success"):
            return {
                "content": [{"type": "text", "text": f"{kind} error: {render_result.get('error')}"}],
                "is_error": True,
            }
        path = render_result["path"]
        meta = render_tool.image_metadata(path)
        b64 = render_tool.png_to_base64(path)
        insert_hint = (
            "When you insert via run_pptxgen_code, use inch values "
            "that preserve the aspect ratio"
            if use_pptxgen else
            "When you insert via run_pptx_code, use Inches() values "
            "that preserve the aspect ratio"
        )
        return {
            "content": [
                {"type": "image",
                 "source": {"type": "base64", "media_type": "image/png", "data": b64}},
                {"type": "text",
                 "text": (
                     f"{kind} saved: {path}\n"
                     f"Dimensions: {meta['width_px']}×{meta['height_px']}px "
                     f"(aspect {meta['aspect']}).\n"
                     f"Inspect the image above BEFORE inserting. If labels are clipped, "
                     f"colors look wrong, or the layout is ugly, regenerate with different "
                     f"parameters. {insert_hint} — e.g. for aspect=2.0 "
                     f"use width=10in / height=5in, not arbitrary dims."
                 )},
            ],
        }

    @tool(
        "make_chart",
        "Render a matplotlib chart to PNG. Use for quantitative data viz (bar/line/pie/scatter/distribution). Code gets `plt` and `np` in scope and `output_path` set. The rendered PNG is returned to you as an image so you can inspect it before inserting.",
        {"code": str, "out_path": str},
    )
    async def make_chart_tool(args):
        r = render_tool.render_matplotlib_chart(args["code"], args["out_path"])
        return _diagram_result(r, "Chart")

    @tool(
        "make_equation",
        "Typeset a math expression to a tightly-cropped TRANSPARENT PNG (matplotlib mathtext: "
        "fractions, integrals, sums, roots, Greek, sub/superscripts, operators — no \\begin "
        "environments). Use for any DISPLAY equation; NEVER write TeX markup as slide text. "
        "latex: the expression WITHOUT surrounding $ (e.g. '\\frac{dy}{dt} = -ky'). color: hex "
        "text color WITH # — MUST match the destination slide's text color (light color on dark "
        "masters, same contrast rule as icons). serif: true on serif-styled decks (Georgia/"
        "Garamond/Palatino body fonts) for Times-like math; false/omit on sans decks. fontsize: "
        "pt at render (44 default; the placed image's size on the slide is what matters — use "
        "the returned aspect to size without distortion). Returns the PNG for inspection before "
        "inserting.",
        {"latex": str, "out_path": str, "color": str, "fontsize": int, "serif": bool},
    )
    async def make_equation_tool(args):
        r = render_tool.render_equation(
            args["latex"], args["out_path"],
            color=args.get("color") or "#1A1D2E",
            fontsize=int(args.get("fontsize") or 44),
            serif=bool(args.get("serif")))
        return _diagram_result(r, "Equation")

    @tool(
        "make_flowchart",
        "Render a Graphviz DOT flowchart to PNG. Legacy — prefer manual PptxGenJS shapes for linear flows, or make_d2_diagram for complex branching diagrams. Supply full DOT source (starting with 'digraph G { ... }').",
        {"dot": str, "out_path": str},
    )
    async def make_flowchart_tool(args):
        r = render_tool.render_graphviz_flowchart(args["dot"], args["out_path"])
        return _diagram_result(r, "Flowchart")

    @tool(
        "make_d2_diagram",
        "Render a D2 diagram (https://d2lang.com) to PNG. Use for complex branching diagrams, architecture diagrams, and decision trees with many interconnected nodes — where automatic edge routing adds real value. For simple linear flows, prefer manual PptxGenJS shapes instead. Use 'dagre' layout for flow, 'elk' for strict hierarchy / right-angle edges, or 'tala' for architecture. Returns the rendered PNG so you can inspect it.",
        {"source": str, "out_path": str, "layout": str},
    )
    async def make_d2_tool(args):
        layout = args.get("layout") or "dagre"
        r = render_tool.render_d2_diagram(args["source"], args["out_path"], layout)
        return _diagram_result(r, "D2 diagram")

    @tool(
        "make_mermaid_diagram",
        "Render a Mermaid diagram (https://mermaid.js.org) to PNG. Best for sequence diagrams, ER diagrams, gantt charts, and state diagrams — Mermaid's sweet spot. For simple linear flows prefer manual PptxGenJS shapes; for complex branching diagrams prefer make_d2_diagram. Supply full Mermaid source (e.g. starts with 'sequenceDiagram' or 'erDiagram').",
        {"source": str, "out_path": str},
    )
    async def make_mermaid_tool(args):
        r = render_tool.render_mermaid_diagram(args["source"], args["out_path"])
        return _diagram_result(r, "Mermaid diagram")

    @tool(
        "make_icon",
        "Render a react-icons icon to PNG. Returns the image so you can inspect "
        "it before inserting into a slide. Get icon_name + icon_pack from find_icon "
        "first — a guessed name often doesn't exist and wastes a turn. icon_pack: "
        "'fa' (Font Awesome), 'md' (Material Design), 'hi2' (Heroicons v2), 'tb' "
        "(Tabler), 'bs' (Bootstrap), 'ai' (Devicons). icon_name: the export name, "
        "e.g. 'FaRocket', 'MdDashboard', 'HiAcademicCap'. color: hex WITHOUT # prefix.",
        {"icon_name": str, "icon_pack": str, "out_path": str, "color": str, "size": int},
    )
    async def make_icon_tool(args):
        if not pptxgen_session:
            return {
                "content": [{"type": "text", "text": "make_icon requires a PptxGenJS session (new build mode)."}],
                "is_error": True,
            }
        result = pptxgen_session.render_icon(
            icon_name=args["icon_name"],
            icon_pack=args.get("icon_pack", "fa"),
            out_path=args["out_path"],
            color=args.get("color", "000000"),
            size=args.get("size", 256),
        )
        if not result.get("success"):
            return {
                "content": [{"type": "text", "text": f"Icon error: {result.get('error')}"}],
                "is_error": True,
            }
        icon_path = result.get("path", args["out_path"])
        b64 = result.get("base64", "")
        content = []
        if b64:
            content.append({
                "type": "image",
                "source": {"type": "base64", "media_type": "image/png", "data": b64},
            })
        content.append({
            "type": "text",
            "text": f"Icon saved: {icon_path} ({result.get('size', 256)}px). "
                    f"Insert via slide.addImage({{ path: '{icon_path}', x: ..., y: ..., w: 0.6, h: 0.6 }}).",
        })
        return {"content": content}

    @tool(
        "find_icon",
        "Search the installed react-icons library for REAL, valid icon names by concept "
        "(e.g. 'handshake', 'lock', 'gauge', 'people'). Returns actual export names and the "
        "pack each lives in — names make_icon is guaranteed to render. ALWAYS use this to get "
        "a name instead of guessing: a guessed icon_name often doesn't exist and wastes a turn.",
        {"query": str},
    )
    async def find_icon_tool(args):
        if not pptxgen_session:
            return {
                "content": [{"type": "text", "text": "find_icon requires a PptxGenJS session (new build mode)."}],
                "is_error": True,
            }
        result = pptxgen_session.find_icon(query=args["query"])
        if not result.get("success"):
            return {
                "content": [{"type": "text", "text": f"find_icon error: {result.get('error')}"}],
                "is_error": True,
            }
        matches = result.get("matches", [])
        if not matches:
            return {"content": [{"type": "text", "text": (
                f"No icons matched '{args['query']}'. Try a simpler or synonym term "
                f"(e.g. 'ban' for no-entry, 'users' for people, 'gauge' for speedometer)."
            )}]}
        lines = "\n".join(f"  {m['name']}  (icon_pack: '{m['pack']}')" for m in matches)
        return {"content": [{"type": "text", "text": (
            f"Real icons matching '{args['query']}' — pass the name + its icon_pack to make_icon:\n{lines}"
        )}]}

    @tool(
        "render_slide",
        "Render one slide of the current .pptx to PNG and return the image so you can visually verify your work. slide_index is 0-based.",
        {"slide_index": int},
    )
    async def render_slide_tool(args):
        render_path = pptx_path
        snap_dir = None
        if use_pptxgen and pptxgen_session:
            snap_dir = tempfile.mkdtemp(prefix="syndara_snap_")
            snap_path = os.path.join(snap_dir, "snap.pptx")
            snap_result = pptxgen_session.snapshot(snap_path)
            if snap_result.get("success") and Path(snap_path).exists():
                render_path = snap_path
                # Persist the latest full deck to the real output path. The
                # snapshot already serialized it, so this is just a file copy —
                # if the Node worker dies before the final save(), the build
                # falls back to this instead of losing everything.
                try:
                    shutil.copyfile(snap_path, pptx_path)
                except Exception:
                    pass

        try:
            b64 = render_tool.render_slide_png_b64(render_path, args["slide_index"])
            return {
                "content": [
                    {
                        "type": "image",
                        "source": {"type": "base64", "media_type": "image/png", "data": b64},
                    },
                    {
                        "type": "text",
                        "text": f"Rendered slide {args['slide_index'] + 1}. Inspect and iterate if anything looks off.",
                    },
                ]
            }
        except Exception as e:
            return {
                "content": [{"type": "text", "text": f"Render failed: {e}"}],
                "is_error": True,
            }
        finally:
            if snap_dir:
                shutil.rmtree(snap_dir, ignore_errors=True)

    @tool(
        "read_summary",
        "Return structural summary of the .pptx. Pass slide_index (0-based) to get detailed info for ONE slide, or omit it for an overview of all slides. For targeted edits, always pass slide_index to get full shape details for the slide you're about to edit.",
        {"slide_index": int},
    )
    async def read_summary_tool(args):
        si = args.get("slide_index")
        data = code_exec_tool.read_pptx_summary(pptx_path, slide_index=si)
        return {"content": [{"type": "text", "text": json.dumps(data)[:8000]}]}

    # ── Choose tools based on engine ────────────────────────────
    if use_pptxgen:
        # New build: PptxGenJS is the primary slide tool
        mcp_tools = [
            run_pptxgen_code_tool,
            make_chart_tool,
            make_equation_tool,
            make_flowchart_tool,
            make_d2_tool,
            make_mermaid_tool,
            make_icon_tool,
            find_icon_tool,
            render_slide_tool,
            read_summary_tool,
        ]
        allowed = [
            "Read", "Write",
            "mcp__pptx__run_pptxgen_code",
            "mcp__pptx__make_chart",
            "mcp__pptx__make_equation",
            "mcp__pptx__make_flowchart",
            "mcp__pptx__make_d2_diagram",
            "mcp__pptx__make_mermaid_diagram",
            "mcp__pptx__make_icon",
            "mcp__pptx__find_icon",
            "mcp__pptx__render_slide",
            "mcp__pptx__read_summary",
        ]
    else:
        # Targeted edit: python-pptx for modifying existing slides
        mcp_tools = [
            run_pptx_code_tool,
            make_chart_tool,
            make_equation_tool,
            make_flowchart_tool,
            make_d2_tool,
            make_mermaid_tool,
            render_slide_tool,
            read_summary_tool,
        ]
        allowed = [
            "Read", "Write",
            "mcp__pptx__run_pptx_code",
            "mcp__pptx__make_chart",
            "mcp__pptx__make_equation",
            "mcp__pptx__make_flowchart",
            "mcp__pptx__make_d2_diagram",
            "mcp__pptx__make_mermaid_diagram",
            "mcp__pptx__render_slide",
            "mcp__pptx__read_summary",
        ]

    pptx_mcp = create_sdk_mcp_server(
        name="pptx",
        version="1.0.0",
        tools=mcp_tools,
    )

    # Shape the system prompt with runtime values
    sys_prompt_template = SYSTEM_PROMPT_PPTXGEN if use_pptxgen else SYSTEM_PROMPT_PPTX
    directive = _DESIGN_DIRECTIVE
    if watermark_info and watermark_info.get("mode"):
        wm_directive = (
            "\n\nWATERMARK RESERVED ZONES (a watermark image is added automatically in post-processing):\n"
            "- EVERY SLIDE: keep content above y=6.5\" in the right third (x > 10\").\n"
        )
        if watermark_info["mode"] == "title_and_bottom_right":
            wm_directive += "- TITLE SLIDE (slide 0): keep y=0 to y=1.1\" clear at the center (x=3.5 to x=9.8).\n"
        wm_directive += "Do NOT add the watermark yourself."
        directive += wm_directive
    if outline.get("plain_backgrounds"):
        visual_motif_text = (
            "VISUAL MOTIF:\n"
            "Use PLAIN backgrounds only — this applies to EVERY slide including\n"
            "title slides. Do NOT add decorative circles, ellipses, blobs, or any\n"
            "addShape background elements on ANY slide. No exceptions for title or\n"
            "section slides. Keep all slides clean: content only, no decorations."
        )
    else:
        visual_motif_text = (
            "VISUAL MOTIF — CRITICAL:\n"
            "The CONTENT master already has a left accent bar. On your first\n"
            "run_pptxgen_code call, ALSO add 1-2 faded accent circles\n"
            "(transparency: 88-92) at consistent positions on every slide. These\n"
            "background circles are what make the deck feel designed. Place them\n"
            "partially off-canvas (negative x/y) for a professional bleed effect.\n"
            "Example for dark title slides:\n"
            "  ts.addShape('ellipse', { x: -1.5, y: 4.0, w: 5, h: 5,\n"
            "      fill: { color: s.accent2 || s.accent, transparency: 85 }, line: { width: 0 } });\n"
            "  ts.addShape('ellipse', { x: 10, y: -1, w: 4, h: 4,\n"
            "      fill: { color: s.bg, transparency: 92 }, line: { width: 0 } });"
        )
    sys_prompt = sys_prompt_template.format(
        outline_path=outline_path,
        pptx_path=pptx_path,
        max_turns=max_turns,
        design_directive=directive,
        max_words_per_slide=outline.get("max_words_per_slide") or 20,
        visual_motif=visual_motif_text,
        # Builds-off courses get a prompt with the reveal-tag teaching removed entirely — the
        # builder can't emit a convention it was never shown. (Extra kwargs are ignored by
        # templates without the placeholder, e.g. the PPTX edit prompt.)
        progressive_builds_section=(
            PROGRESSIVE_BUILDS_SECTION if outline.get("progressive_builds", True) else ""),
    )
    from .base import STYLE_RULE
    sys_prompt += STYLE_RULE

    # Per-run API-key override: if a build was started with an owner's alternate Anthropic key,
    # forward it (plus the rest of the process env) to the Claude Code subprocess. When there's no
    # override we leave env unset so the subprocess uses its default environment as before.
    _opts_kwargs: dict = {}
    try:
        from .. import keyring as _keyring
        _ovr = _keyring.anthropic_override()
        if _ovr:
            import os as _os
            _opts_kwargs["env"] = {**_os.environ, "ANTHROPIC_API_KEY": _ovr}
    except Exception:
        pass

    import os as _os_env
    options = ClaudeAgentOptions(
        model=BUILDER_MODEL,
        mcp_servers={"pptx": pptx_mcp},
        allowed_tools=allowed,
        permission_mode="acceptEdits",
        system_prompt=sys_prompt,
        max_turns=max_turns,
        # The SDK's default 1MB per-JSON-message buffer killed a whole module build when one
        # render_slide preview PNG base64'd past it (job 203 m5: "JSON message exceeded maximum
        # buffer size"). 8MB is far above any legitimate tool result while still bounding memory;
        # render_tool also downscales oversized previews so messages stay well under it.
        max_buffer_size=int(_os_env.environ.get("SYNDARA_SDK_MAX_BUFFER", str(8 * 1024 * 1024))),
        **_opts_kwargs,
    )

    # ── Drive the agent loop ─────────────────────────────────────
    import time as _time
    _t0 = _time.time()
    _tool_count = 0
    _label = f"Builder.{outline.get('module_id', '?')}"
    _mode = "targeted-edit" if is_targeted_edit else "full-build"
    print(f"[{_label}] START mode={_mode} max_turns={max_turns}", flush=True)

    # Wall-clock deadline: this loop previously had NO overall timeout, so a wedged SDK/CLI
    # session hung the module forever while holding a concurrency slot. On expiry we break and
    # fall through to the same salvage path as a non-success end.
    _deadline_s = int(os.environ.get("SYNDARA_BUILDER_DEADLINE", "3600"))
    _nonsuccess: Optional[str] = None
    try:
        final_result: Optional[str] = None
        _tool_names: dict = {}       # tool name -> count
        _action_counts: dict = {}    # "name(input-preview)" -> count, to spot repeated/spinning actions
        async for message in query(prompt=prompt, options=options):
            if _time.time() - _t0 > _deadline_s:
                _nonsuccess = "deadline_exceeded"
                print(f"[{_label}] DEADLINE ({_deadline_s}s) exceeded — stopping the session and "
                      f"salvaging the persisted deck", flush=True)
                break
            if isinstance(message, AssistantMessage):
                from claude_agent_sdk import ToolUseBlock
                for block in message.content:
                    if isinstance(block, ToolUseBlock):
                        _tool_count += 1
                        elapsed = _time.time() - _t0
                        inp_preview = json.dumps(block.input)[:120] if block.input else ""
                        _tool_names[block.name] = _tool_names.get(block.name, 0) + 1
                        _ak = f"{block.name}({inp_preview})"
                        _action_counts[_ak] = _action_counts.get(_ak, 0) + 1
                        print(
                            f"[{_label}] tool {_tool_count}/{max_turns} "
                            f"@ {elapsed:.0f}s — {block.name}({inp_preview})",
                            flush=True,
                        )
            elif isinstance(message, ResultMessage):
                elapsed = _time.time() - _t0
                final_result = getattr(message, "result", None) or getattr(message, "content", None) or ""
                cost = getattr(message, "total_cost_usd", None)
                usage = getattr(message, "usage", None) or {}
                try:
                    from .base import report_exact_cost
                    report_exact_cost(_label, cost)   # exact, cache-aware; per-run sink
                except Exception:
                    pass
                # Structured build summary for the generation trace: surfaces spinning
                # (one action repeated many times) and tool-budget exhaustion.
                try:
                    from .base import report_build_summary
                    _worst = max(_action_counts.items(), key=lambda kv: kv[1], default=(None, 0))
                    report_build_summary(_label, {
                        "turns": getattr(message, "num_turns", None),
                        "tools": _tool_count,
                        "tools_by_name": _tool_names,
                        "elapsed_s": round(elapsed, 1),
                        "cost_usd": cost,
                        "stop": message.stop_reason,
                        "status": message.subtype,
                        "hit_max_turns": bool(_tool_count >= max_turns or message.subtype == "error_max_turns"),
                        "repeated_actions": sum(1 for v in _action_counts.values() if v >= 3),
                        "worst_repeat": ({"action": _worst[0], "count": _worst[1]} if _worst[1] >= 3 else None),
                    })
                except Exception:
                    pass
                print(
                    f"[{_label}] DONE {elapsed:.0f}s · turns={message.num_turns} "
                    f"· tools={_tool_count} · cost=${cost or '?'} "
                    f"· in={usage.get('input_tokens', '?')} out={usage.get('output_tokens', '?')} "
                    f"· stop={message.stop_reason} · status={message.subtype}",
                    flush=True,
                )
                if message.subtype != "success":
                    # Do NOT raise here: render_slide snapshots the deck to pptx_path on every
                    # self-check, so a max-turns/aborted session usually leaves a complete or
                    # near-complete valid deck on disk. Raising threw that paid work away (two
                    # prior incidents dropped whole 40-slide modules). Record it and let the
                    # salvage check after the loop decide.
                    _nonsuccess = message.subtype
                    print(f"[{_label}] NON-SUCCESS end ({message.subtype}) — will try to salvage "
                          f"the persisted deck", flush=True)
    except Exception as e:
        traceback.print_exc()
        # Try to save whatever we have before raising
        if pptxgen_session:
            try:
                pptxgen_session.save()
            except Exception:
                pass
        raise RuntimeError(f"Claude Code builder failed: {e}") from e

    # ── Finalize the PPTX file (PptxGenJS builds) ──────────────
    if pptxgen_session:
        try:
            save_result = pptxgen_session.save()
            if not save_result.get("success"):
                print(f"[{_label}] PPTXGEN SAVE WARNING: {save_result.get('error')}", flush=True)
                if Path(pptx_path).exists():
                    print(
                        f"[{_label}] RECOVERED: final save failed (worker died), "
                        f"falling back to last persisted snapshot at {pptx_path}",
                        flush=True,
                    )
            else:
                print(
                    f"[{_label}] PPTXGEN SAVED: {save_result.get('slide_count')} slides to {pptx_path}",
                    flush=True,
                )
        except Exception as e:
            print(f"[{_label}] PPTXGEN SAVE ERROR: {e}", flush=True)

    if not Path(pptx_path).exists():
        raise RuntimeError(
            f"Claude Code builder returned but no pptx at {pptx_path}. The PptxGenJS "
            f"worker died before any deck was persisted (no render snapshot taken). "
            f"Final message: {(final_result or '')[:400]}"
        )

    if _nonsuccess:
        # Salvage decision for a session that ended non-success (error_max_turns, CLI abort,
        # deadline): accept the persisted deck when it's a valid ZIP holding a plausible share
        # of the planned slides — Visual QA and creator review catch the gaps. Only raise when
        # nothing usable exists; the old behavior threw away complete 40-slide decks.
        import re as _re
        import zipfile as _zf
        _ok = False
        _n = 0
        try:
            with _zf.ZipFile(pptx_path) as _z:
                _n = sum(1 for nm in _z.namelist() if _re.match(r"ppt/slides/slide\d+\.xml$", nm))
            _need = max(3, int(target_slides * 0.6)) if target_slides else 3
            _ok = _n >= _need
            print(f"[{_label}] SALVAGE CHECK after {_nonsuccess}: {_n} slides on disk, "
                  f"need >= {_need} -> {'ACCEPT (degraded)' if _ok else 'reject'}", flush=True)
        except Exception as _se:
            print(f"[{_label}] salvage check failed: {_se}", flush=True)
        if not _ok:
            raise RuntimeError(
                f"Claude Code builder ended non-success: {_nonsuccess}, and no salvageable "
                f"deck was persisted ({_n} slides). {(final_result or '')[:400]}")

    # Strip OPC-violating directory entries that PptxGenJS/JSZip leaves in the
    # chart data workbooks — otherwise PowerPoint shows a 'repair' dialog.
    try:
        from ..tools.pptx_tool import sanitize_pptx
        if sanitize_pptx(pptx_path):
            print(f"[{_label}] sanitized OOXML (removed JSZip directory entries)", flush=True)
    except Exception as _e:
        print(f"[{_label}] sanitize step skipped: {_e}", flush=True)

    # Post-build validation: check that slides have images
    try:
        from pptx import Presentation as _Prs
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        _prs = _Prs(pptx_path)
        total = len(_prs.slides)
        with_images = sum(
            1 for sl in _prs.slides
            if any(
                getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
                or getattr(sh, "image", None) is not None
                for sh in sl.shapes
            )
        )
        print(
            f"[{_label}] POST-BUILD CHECK: {with_images}/{total} slides have images "
            f"({'OK' if with_images > total * 0.5 else 'LOW — visuals may be missing'})",
            flush=True,
        )
    except Exception as _e:
        print(f"[{_label}] POST-BUILD CHECK failed: {_e}", flush=True)

    return pptx_path


# ── Sync wrapper so web_runner can run_in_thread like the other builders ─
def build(
    outline: dict,
    output_dir: str,
    style: str = "syndara",
    reviewer_feedback: Optional[dict] = None,
    watermark_info: Optional[dict] = None,
) -> str:
    """Synchronous-looking entry point. Spins its own asyncio loop so this
    is safe to call from asyncio.to_thread().

    Corrupt-output safety net: the agentic loop can (rarely) leave PptxGenJS in a bad state so the
    saved file isn't a valid .pptx/ZIP while save() still reports success (e.g. a destructive slide
    reset). That's a silent way to kill a paid build, so validate the output here and rebuild ONCE
    from a fresh session before giving up. Only fires on an already-broken build, so no cost in the
    normal case. Centralized here so every caller (single-deck, course, revamp) is protected."""
    import asyncio
    import zipfile as _zip
    _label = outline.get("module_id", "?")
    path = ""
    for attempt in range(2):   # original + one retry
        path = asyncio.run(
            build_slides_with_claude_code(outline, output_dir, style, reviewer_feedback,
                                          watermark_info=watermark_info)
        )
        if path and Path(path).exists() and _zip.is_zipfile(path):
            return path
        if attempt == 0:
            print(f"[Builder.{_label}] output is not a valid pptx (corrupt/non-ZIP) — "
                  f"rebuilding ONCE from a fresh session", flush=True)
    # Still bad after the retry — return the path and let the caller's validation raise a clear error.
    print(f"[Builder.{_label}] output still invalid after retry — surfacing to caller", flush=True)
    return path


class ClaudeCodeSlideBuilder:
    """Thin class for parity with AgenticSlideBuilder's interface."""

    def build(
        self,
        outline: dict,
        output_dir: str,
        style: str = "syndara",
        reviewer_feedback: Optional[dict] = None,
        watermark_info: Optional[dict] = None,
    ) -> str:
        return build(outline, output_dir, style, reviewer_feedback, watermark_info=watermark_info)
