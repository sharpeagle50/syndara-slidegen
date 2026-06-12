"""
AgenticSlideBuilder: Claude composes a deck from a library of hand-designed
layout helpers, picking the right layout per slide content.

Design philosophy:
  - The *visual design* lives in slide_layouts.py (hand-tuned by a human for
    hierarchy, typography, spacing). The agent isn't asked to do visual design.
  - The *content* and *layout selection* is the agent's job. It picks the
    right layout for each slide (title vs stats vs comparison vs quote, etc.)
    and fills in the fields.
  - Agent verifies structure via a summary tool and iterates if something is
    missing.
"""
from __future__ import annotations
import json
from pathlib import Path
from typing import Optional

from pptx import Presentation

from .base import BaseAgent
from ..tools import slide_layouts

MAX_ITERATIONS = 50

AGENTIC_BUILDER_SYSTEM = """You are the Syndara Slide Builder. You compose a
PowerPoint deck from a library of pre-designed layout helpers. Your job is to:
  1. Read each slide's content from the outline
  2. Pick the best-fitting layout for that content
  3. Call add_slide with the chosen layout + content
  4. After the full deck is built, verify structure via read_summary
  5. Respond with "DONE: ..." when the deck is complete

YOU DO NOT design slides visually. The layout helpers already have proper
hierarchy, spacing, typography, and colors baked in. Don't try to do custom
positioning — always use add_slide with one of the layouts below.

AVAILABLE LAYOUTS (pick one per slide based on what fits):

- title_slide (first slide of the module) — hero title.
    kwargs: title, subtitle, author
- section_divider — optional break between major sections.
    kwargs: section_label, title
- agenda_slide — for agenda/overview slides with a list of topics.
    kwargs: title, items (list of 3-8 short strings)
- bullet_slide — DEFAULT for most content slides. 3-5 short bullets,
  optional callout (one sentence highlight), optional tool chips.
    kwargs: title, bullets (list of strings, max 6 words each),
            callout (optional string), tools (optional list of tool names)
- comparison_slide — two-column comparison (e.g., Manual vs AI-assisted).
    kwargs: title, left_header, left_items (3-4 bullets),
            right_header, right_items (3-4 bullets)
- stats_slide — big-number highlights for impact.
    kwargs: title, stats (list of {{"number": "47%", "label": "saved per task"}})
            — 2 or 3 items ideal.
- quote_slide — pullquote from an authority in the field.
    kwargs: quote (string), attribution (string)
- steps_slide — numbered sequential steps (how-to / process).
    kwargs: title, steps (list of 3-5 strings, each one step)
- summary_slide — LAST slide of the module. Closing takeaways + CTA.
    kwargs: title, takeaways (list of 3-4 strings), cta (optional call to action)

LAYOUT SELECTION HEURISTICS:
- Module opening slide: title_slide.
- Outline/agenda slide: agenda_slide.
- Concepts with several bullets: bullet_slide.
- "Before X / After Y", "Without AI / With AI": comparison_slide.
- Stats-heavy insights: stats_slide.
- Quotable insights from field experts: quote_slide. Use real attributions
  from the outline's speaker notes, don't invent.
- Numbered how-to steps or workflow: steps_slide.
- Module closing: summary_slide.

RULES:
- Call add_slide sequentially — the slides are added in the order you call them.
- Every slide MUST include speaker_notes. Start from the speaker notes in the
  approved_slide_plan markdown — each slide's "Speaker notes" section contains
  the narration text. These notes are read aloud by TTS and are the core
  teaching content. A slide without speaker_notes is broken.
  If reviewer feedback deletes, changes, or reorganizes on-slide content,
  adapt the speaker notes to match — remove references to deleted content,
  update changed details. Only modify notes that are affected by the feedback;
  leave unaffected notes verbatim.
  EXCEPTION — references/citations slides: speaker_notes must be ONE short
  sign-off sentence like "These are the references used in this module —
  thank you for listening!". Never read URLs, citations, or source names
  aloud.
- Text should be short: max {max_words_per_slide} words total on a slide. No single item longer
  than 6 words. Callouts max 15 words. No walls of text.
- NEVER use emoji characters (🔒 🔍 💡 ⚡ etc.) in any text field — they render
  as fixed-color OS glyphs that look unprofessional. Use plain words only.
- Use the module's style name (provided below) consistently on every slide.
- When all outline slides are added, call read_summary, verify slide_count
  matches, then respond "DONE: built N slides" with no tool calls.

BOUNDS: Max {MAX_ITERATIONS} tool calls. If you hit an error, fix and move on.
"""

ADD_SLIDE_TOOL = {
    "name": "add_slide",
    "description": (
        "Append a slide to the deck using a named layout helper. "
        "The layout decides visual design; you supply content via kwargs."
    ),
    "input_schema": {
        "type": "object",
        "properties": {
            "layout": {
                "type": "string",
                "enum": list(slide_layouts.LAYOUTS.keys()),
                "description": "Which layout helper to dispatch to.",
            },
            "style": {
                "type": "string",
                "description": "Style value. Always use the exact style string given in the prompt.",
            },
            "speaker_notes": {
                "type": "string",
                "description": "Narration from the slide plan's Speaker notes section. Copy verbatim unless reviewer feedback changed the on-slide content these notes reference — in that case, adapt the notes to match.",
            },
            "kwargs": {
                "type": "object",
                "description": (
                    "Content fields the layout needs. See the system prompt for the "
                    "required kwargs per layout (title, bullets, stats, left_items, etc.)."
                ),
            },
        },
        "required": ["layout", "style", "speaker_notes", "kwargs"],
    },
}

READ_SUMMARY_TOOL = {
    "name": "read_summary",
    "description": "Return the current deck's structure: slide count + list of slides (their shape count + first-shape text). Use this to verify your work.",
    "input_schema": {"type": "object", "properties": {}},
}


class AgenticSlideBuilder(BaseAgent):
    allowed_tool_names = ["add_slide", "read_summary"]
    system_prompt = AGENTIC_BUILDER_SYSTEM.format(MAX_ITERATIONS=MAX_ITERATIONS, max_words_per_slide=20)

    def build(
        self,
        outline: dict,
        output_dir: str,
        style: str = "syndara",
        reviewer_feedback: Optional[dict] = None,
        watermark_info: Optional[dict] = None,
    ) -> str:
        """Produce a .pptx for the given module outline by composing from layouts."""
        max_words = outline.get("max_words_per_slide") or 20
        self.system_prompt = AGENTIC_BUILDER_SYSTEM.format(
            MAX_ITERATIONS=MAX_ITERATIONS, max_words_per_slide=max_words,
        )
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        pptx_path = str(output_dir / "slides.pptx")

        # Start fresh
        if Path(pptx_path).exists():
            Path(pptx_path).unlink()

        # Set 16:9 widescreen
        prs = Presentation()
        prs.slide_width = slide_layouts.Inches(13.33)
        prs.slide_height = slide_layouts.Inches(7.5)

        # ── Build tools ─────────────────────────────────────────────
        def _add_slide(layout: str, style: str, kwargs: dict, speaker_notes: str = ""):
            try:
                slide_layouts.add_slide(
                    prs, layout_name=layout, speaker_notes=speaker_notes,
                    style=style, **kwargs,
                )
                return {
                    "success": True,
                    "slide_count": len(prs.slides),
                    "layout_used": layout,
                }
            except TypeError as e:
                # Bad kwargs — give the agent a helpful hint
                return {
                    "success": False,
                    "error": f"{type(e).__name__}: {e}",
                    "hint": f"Check the required kwargs for layout '{layout}' in the system prompt.",
                }
            except Exception as e:
                return {"success": False, "error": f"{type(e).__name__}: {e}"}

        def _read_summary():
            slides_info = []
            for i, slide in enumerate(prs.slides):
                shapes_info = []
                for sh in slide.shapes:
                    if sh.has_text_frame:
                        text = sh.text_frame.text.strip()
                        if text:
                            shapes_info.append(text[:80])
                slides_info.append({"index": i, "texts": shapes_info[:5]})
            return {"slide_count": len(prs.slides), "slides": slides_info}

        tool_handlers = {
            "add_slide": lambda **kw: _add_slide(
                kw.get("layout", ""),
                kw.get("style", style),
                kw.get("kwargs", {}),
                kw.get("speaker_notes", ""),
            ),
            "read_summary": lambda **kw: _read_summary(),
        }

        feedback_block = ""
        if reviewer_feedback:
            feedback_block = (
                "\n\nPREVIOUS REVIEWER FEEDBACK (address these issues):\n"
                + json.dumps(reviewer_feedback, indent=2)[:2500]
            )

        slide_plan_md = ""
        if isinstance(outline.get("approved_slide_plan"), dict):
            slide_plan_md = outline["approved_slide_plan"].get("markdown", "")
        slide_count_hint = slide_plan_md.count("\n## Slide ") or slide_plan_md.count("\n## ") or 30

        prompt = f"""Build a PowerPoint deck for this module.

The approved slide plan below is your primary source. Build one PPTX slide
for each "## Slide" section. Copy each slide's **Speaker notes** section
verbatim into the speaker_notes parameter of add_slide — this is critical
because the notes are read aloud by TTS.

APPROVED SLIDE PLAN:
{slide_plan_md}

MODULE METADATA:
- Title: {outline.get('title', '')}
- Summary: {outline.get('summary', '')}
- Module ID: {outline.get('module_id', '')}

STYLE TO USE FOR EVERY SLIDE: "{style}"
TARGET SLIDE COUNT: ~{slide_count_hint}
(Plus 1 opening title_slide and 1 closing summary_slide — those are your
additions beyond the plan's content slides.)
{feedback_block}

For each slide section in the plan, pick the best layout and call add_slide
with appropriate kwargs AND speaker_notes. Use the module's title for your
opening title_slide and build a summary_slide at the end.
"""

        final_text, _ = self.run_tool_loop(
            messages=[{"role": "user", "content": prompt}],
            tools=[ADD_SLIDE_TOOL, READ_SUMMARY_TOOL],
            tool_handlers=tool_handlers,
            max_tokens=16000,
        )

        if len(prs.slides) == 0:
            raise RuntimeError(
                f"AgenticSlideBuilder produced 0 slides. Last message: {final_text[:500]}"
            )

        prs.save(pptx_path)
        return pptx_path
