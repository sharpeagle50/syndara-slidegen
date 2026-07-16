"""Post-processing watermark insertion into PPTX slides."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image


SLIDE_W = Emu(Inches(13.33))
SLIDE_H = Emu(Inches(7.5))

BR_MAX_W = Emu(Inches(2.5))
BR_MAX_H = Emu(Inches(0.5))
BR_MARGIN_RIGHT = Emu(Inches(0.5))
BR_MARGIN_BOTTOM = Emu(Inches(0.3))

CT_MAX_W = Emu(Inches(6.0))
CT_MAX_H = Emu(Inches(0.7))
CT_Y = Emu(Inches(0.3))

# Creator-chosen size multiplier bounds (1.0 = default box). strip_watermark must match up to the max
# so a scaled-up watermark is still recognized + removed on re-apply (else re-watermark duplicates it).
_MIN_SCALE = 0.25
_MAX_SCALE = 3.0


def _image_size(image_path: str) -> tuple[int, int]:
    with Image.open(image_path) as img:
        return img.size


def _fit(img_w: int, img_h: int, max_w: int, max_h: int) -> tuple[Emu, Emu]:
    ratio_w = max_w / img_w
    ratio_h = max_h / img_h
    scale = min(ratio_w, ratio_h)
    return Emu(int(img_w * scale)), Emu(int(img_h * scale))


_STRIP_TOL = Emu(Inches(0.08))   # position-match tolerance when identifying a watermark picture


def strip_watermark(pptx_path: str) -> int:
    """Remove watermark pictures previously added by apply_watermark, identified by their fixed
    anchor: the bottom-right corner (right/bottom edges flush to the slide with the known margins),
    plus the center-top title stamp. Returns how many pictures were removed. Best-effort per shape —
    a picture that doesn't match the anchor (i.e. real content) is left untouched."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def _near(a, b) -> bool:
        return abs(int(a) - int(b)) <= int(_STRIP_TOL)

    prs = Presentation(pptx_path)
    br_right = int(SLIDE_W) - int(BR_MARGIN_RIGHT)      # where the BR watermark's right edge sits
    br_bottom = int(SLIDE_H) - int(BR_MARGIN_BOTTOM)    # ...and its bottom edge
    removed = 0
    for slide in prs.slides:
        for shp in list(slide.shapes):
            try:
                if shp.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                left, top, w, h = int(shp.left), int(shp.top), int(shp.width), int(shp.height)
                # bottom-right stamp: right+bottom edges flush to the slide, within the BR size cap
                is_br = (_near(left + w, br_right) and _near(top + h, br_bottom)
                         and w <= int(BR_MAX_W * _MAX_SCALE) + int(_STRIP_TOL) and h <= int(BR_MAX_H * _MAX_SCALE) + int(_STRIP_TOL))
                # center-top title stamp: top at CT_Y, horizontally centered, within the CT size cap
                is_ct = (_near(top, CT_Y) and _near(left + w // 2, int(SLIDE_W) // 2)
                         and w <= int(CT_MAX_W * _MAX_SCALE) + int(_STRIP_TOL) and h <= int(CT_MAX_H * _MAX_SCALE) + int(_STRIP_TOL))
                if is_br or is_ct:
                    shp._element.getparent().remove(shp._element)
                    removed += 1
            except Exception:
                continue
    if removed:
        prs.save(pptx_path)
        from .pptx_tool import sanitize_pptx
        sanitize_pptx(pptx_path)
    return removed


def apply_watermark(pptx_path: str, watermark_image_path: str, mode: str, scale: float = 1.0):
    """scale multiplies the watermark's max box (1.0 = default). The bottom-right anchor (right/bottom
    margins) is scale-invariant, so strip_watermark still matches a resized watermark for re-apply."""
    if mode not in ("bottom_right", "title_and_bottom_right"):
        return
    try:
        scale = max(_MIN_SCALE, min(_MAX_SCALE, float(scale or 1.0)))
    except (TypeError, ValueError):
        scale = 1.0

    prs = Presentation(pptx_path)
    img_w, img_h = _image_size(watermark_image_path)

    br_w, br_h = _fit(img_w, img_h, int(BR_MAX_W * scale), int(BR_MAX_H * scale))
    br_left = Emu(SLIDE_W - BR_MARGIN_RIGHT - br_w)
    br_top = Emu(SLIDE_H - BR_MARGIN_BOTTOM - br_h)

    for i, slide in enumerate(prs.slides):
        if mode == "title_and_bottom_right" and i == 0:
            ct_w, ct_h = _fit(img_w, img_h, int(CT_MAX_W * scale), int(CT_MAX_H * scale))
            ct_left = Emu((SLIDE_W - ct_w) // 2)
            slide.shapes.add_picture(watermark_image_path, ct_left, CT_Y, ct_w, ct_h)
        else:
            slide.shapes.add_picture(watermark_image_path, br_left, br_top, br_w, br_h)

    prs.save(pptx_path)

    # python-pptx's re-save runs after the builder's own sanitize pass and can
    # leave OPC/ZIP quirks that trigger PowerPoint's spurious "repair" dialog —
    # clean the file we just wrote so a watermarked deck always opens cleanly.
    from .pptx_tool import sanitize_pptx
    sanitize_pptx(pptx_path)
