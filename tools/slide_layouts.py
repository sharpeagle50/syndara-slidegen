"""
Hand-designed slide layouts. Each function builds one well-composed slide
with proper visual hierarchy, spacing, and typography. The agent picks the
right layout per slide content.

A slide is 13.33 × 7.5 inches (16:9 widescreen).
"""
from __future__ import annotations
import json
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── Style palette ───────────────────────────────────────────────────

STYLE_PALETTES = {
    "syndara": {
        "bg": "#F7F8FC", "surface": "#FFFFFF", "accent": "#4361EE",
        "accent2": "#7C3AED", "highlight": "#EEF2FF", "border": "#C7D2FE",
        "text": "#1A1D2E", "subtext": "#6B7280",
        "title_font": "Calibri", "body_font": "Calibri",
    },
    "professional": {
        "bg": "#FFFFFF", "surface": "#F8FAFC", "accent": "#1E40AF",
        "accent2": "#0369A1", "highlight": "#EFF6FF", "border": "#BFDBFE",
        "text": "#0F172A", "subtext": "#475569",
        "title_font": "Georgia", "body_font": "Calibri",
    },
    "academic": {
        "bg": "#FDFCF7", "surface": "#F4EFE1", "accent": "#7C2D12",
        "accent2": "#92400E", "highlight": "#FEF3C7", "border": "#D6D3D1",
        "text": "#1C1917", "subtext": "#57534E",
        "title_font": "Georgia", "body_font": "Georgia",
    },
    "friendly": {
        "bg": "#FEF3C7", "surface": "#FDE68A", "accent": "#D97706",
        "accent2": "#EA580C", "highlight": "#FFF7ED", "border": "#FDBA74",
        "text": "#1F2937", "subtext": "#6B7280",
        "title_font": "Calibri", "body_font": "Calibri",
    },
    "bold": {
        "bg": "#18181B", "surface": "#27272A", "accent": "#F472B6",
        "accent2": "#A78BFA", "highlight": "#3F3F46", "border": "#52525B",
        "text": "#FAFAFA", "subtext": "#A1A1AA",
        "title_font": "Calibri", "body_font": "Calibri",
    },
    "midnight": {
        "bg": "#0F172A", "surface": "#1E293B", "accent": "#38BDF8",
        "accent2": "#818CF8", "highlight": "#1E3A5F", "border": "#334155",
        "text": "#F1F5F9", "subtext": "#94A3B8",
        "title_font": "Georgia", "body_font": "Calibri",
    },
    "forest": {
        "bg": "#F0FDF4", "surface": "#DCFCE7", "accent": "#16A34A",
        "accent2": "#0D9488", "highlight": "#ECFDF5", "border": "#86EFAC",
        "text": "#14532D", "subtext": "#4B5563",
        "title_font": "Georgia", "body_font": "Georgia",
    },
    "coral": {
        "bg": "#FFF5F5", "surface": "#FEE2E2", "accent": "#F43F5E",
        "accent2": "#FB923C", "highlight": "#FFF1F2", "border": "#FECACA",
        "text": "#1C1917", "subtext": "#6B7280",
        "title_font": "Calibri", "body_font": "Calibri",
    },
    "slate": {
        "bg": "#F8FAFC", "surface": "#F1F5F9", "accent": "#475569",
        "accent2": "#64748B", "highlight": "#E2E8F0", "border": "#CBD5E1",
        "text": "#0F172A", "subtext": "#64748B",
        "title_font": "Calibri", "body_font": "Calibri",
    },
    "plum": {
        "bg": "#FAF5FF", "surface": "#F3E8FF", "accent": "#9333EA",
        "accent2": "#C026D3", "highlight": "#F5F3FF", "border": "#D8B4FE",
        "text": "#1E1B4B", "subtext": "#6B7280",
        "title_font": "Georgia", "body_font": "Calibri",
    },
    # ── Anthropic-curated palettes (proven in Claude Code PPTX skill) ──
    "midnight_executive": {
        "bg": "#1E2761", "surface": "#2A3370", "accent": "#CADCFC",
        "accent2": "#7EC8E3", "highlight": "#283575", "border": "#3D4F8F",
        "text": "#FFFFFF", "subtext": "#B8C5E8",
        "title_font": "Georgia", "body_font": "Calibri",
    },
    "forest_moss": {
        "bg": "#F5F5F5", "surface": "#FFFFFF", "accent": "#2C5F2D",
        "accent2": "#97BC62", "highlight": "#EDF5ED", "border": "#C5D9C6",
        "text": "#1A2E1A", "subtext": "#5A6B5A",
        "title_font": "Georgia", "body_font": "Calibri",
    },
    "coral_energy": {
        "bg": "#FFFAF9", "surface": "#FFFFFF", "accent": "#F96167",
        "accent2": "#2F3C7E", "highlight": "#FEF0E5", "border": "#FCCACA",
        "text": "#2F3C7E", "subtext": "#7A7E9E",
        "title_font": "Arial Black", "body_font": "Arial",
    },
    "warm_terracotta": {
        "bg": "#E7E8D1", "surface": "#F2F0E4", "accent": "#B85042",
        "accent2": "#A7BEAE", "highlight": "#F0EDDF", "border": "#C9C4A8",
        "text": "#3D2B24", "subtext": "#7A6E5E",
        "title_font": "Palatino", "body_font": "Garamond",
    },
    "ocean_gradient": {
        "bg": "#21295C", "surface": "#2A3470", "accent": "#1C7293",
        "accent2": "#065A82", "highlight": "#283575", "border": "#3D5A8F",
        "text": "#FFFFFF", "subtext": "#A8C4D8",
        "title_font": "Trebuchet MS", "body_font": "Calibri",
    },
    "charcoal_minimal": {
        "bg": "#F2F2F2", "surface": "#FFFFFF", "accent": "#36454F",
        "accent2": "#5E7283", "highlight": "#E8EAED", "border": "#CBD0D5",
        "text": "#212121", "subtext": "#6B7280",
        "title_font": "Calibri", "body_font": "Calibri Light",
    },
    "teal_trust": {
        "bg": "#F0FAFA", "surface": "#FFFFFF", "accent": "#028090",
        "accent2": "#00A896", "highlight": "#E0F5F2", "border": "#90D5CC",
        "text": "#1A2E2E", "subtext": "#5A7A7A",
        "title_font": "Trebuchet MS", "body_font": "Calibri",
    },
    "berry_cream": {
        "bg": "#ECE2D0", "surface": "#F5EFE3", "accent": "#6D2E46",
        "accent2": "#A26769", "highlight": "#F0E8DB", "border": "#D4C4AC",
        "text": "#2E1A22", "subtext": "#7A6A5E",
        "title_font": "Palatino", "body_font": "Garamond",
    },
    "sage_calm": {
        "bg": "#F4F7F5", "surface": "#FFFFFF", "accent": "#50808E",
        "accent2": "#69A297", "highlight": "#E8F0EC", "border": "#B4CEC5",
        "text": "#1E3330", "subtext": "#5E7A72",
        "title_font": "Georgia", "body_font": "Calibri",
    },
}


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


_COLOR_KEYS = {"bg", "surface", "accent", "accent2", "highlight", "border", "text", "subtext"}


def get_palette(style: str) -> dict:
    """Get palette dict (hex strings) by preset name or JSON custom palette."""
    if style in STYLE_PALETTES:
        return STYLE_PALETTES[style]
    try:
        custom = json.loads(style)
        if isinstance(custom, dict) and "bg" in custom and "accent" in custom:
            base = dict(STYLE_PALETTES["professional"])
            base.update(custom)
            return base
    except (json.JSONDecodeError, TypeError):
        pass
    return STYLE_PALETTES["professional"]


def _palette(style: str) -> dict:
    p = get_palette(style)
    return {k: (_rgb(v) if k in _COLOR_KEYS else v)
            for k, v in p.items()}


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _tb(slide, text, left, top, width, height, *, size=18, bold=False,
        color=None, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    if color is not None:
        run.font.color.rgb = color
    return box


def _set_notes(slide, notes: str):
    if not notes:
        return
    slide.notes_slide.notes_text_frame.text = notes


def _blank_slide(prs, pal):
    layout = prs.slide_layouts[6]  # "Blank"
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, pal["bg"])
    return slide


def _accent_bar(slide, pal, left, top, width=0.08, height=0.5):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = pal["accent"]
    bar.line.fill.background()
    return bar


# ── Layout helpers — each returns the created slide ──────────────────

def title_slide(prs, title: str, subtitle: str = "", author: str = "", style: str = "syndara"):
    """Hero title slide — centered title + optional subtitle + author line."""
    pal = _palette(style)
    slide = _blank_slide(prs, pal)
    _tb(slide, title, 1.0, 2.8, 11.33, 1.5,
        size=44, bold=True, color=pal["text"], align=PP_ALIGN.CENTER, font=pal["title_font"])
    if subtitle:
        _tb(slide, subtitle, 1.0, 4.3, 11.33, 0.7,
            size=20, color=pal["accent"], align=PP_ALIGN.CENTER, font=pal["body_font"])
    if author:
        _tb(slide, author, 1.0, 6.5, 11.33, 0.4,
            size=14, color=pal["subtext"], align=PP_ALIGN.CENTER, font=pal["body_font"])
    return slide


def section_divider(prs, section_label: str, title: str, style: str = "syndara"):
    """Section break — small label above big title, with accent bar."""
    pal = _palette(style)
    slide = _blank_slide(prs, pal)
    _tb(slide, section_label.upper(), 1.0, 3.0, 11.33, 0.4,
        size=14, color=pal["accent"], align=PP_ALIGN.CENTER, font=pal["body_font"])
    _tb(slide, title, 1.0, 3.5, 11.33, 1.5,
        size=40, bold=True, color=pal["text"], align=PP_ALIGN.CENTER, font=pal["title_font"])
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.16), Inches(5.1), Inches(1.0), Inches(0.1),
    )
    bar.fill.solid(); bar.fill.fore_color.rgb = pal["accent"]; bar.line.fill.background()
    return slide


def agenda_slide(prs, title: str, items: list[str], style: str = "syndara"):
    """TOC / course overview — numbered list of modules or topics."""
    pal = _palette(style)
    slide = _blank_slide(prs, pal)
    _accent_bar(slide, pal, 0.8, 0.7, 0.08, 0.5)
    _tb(slide, title, 1.0, 0.6, 11.33, 0.7,
        size=32, bold=True, color=pal["text"], font=pal["title_font"])
    for i, item in enumerate(items[:8]):
        _tb(slide, str(i + 1).zfill(2), 1.0, 1.8 + i * 0.65, 0.7, 0.55,
            size=22, color=pal["accent"], bold=True, font=pal["title_font"])
        _tb(slide, item, 1.8, 1.8 + i * 0.65, 11.0, 0.55,
            size=18, color=pal["text"], font=pal["body_font"])
    return slide


def bullet_slide(prs, title: str, bullets: list[str], callout: str = "",
                 tools: Optional[list[str]] = None, style: str = "syndara"):
    """Standard content slide: title + 3-5 bullets + optional callout + tool chips."""
    pal = _palette(style)
    slide = _blank_slide(prs, pal)
    _accent_bar(slide, pal, 0.8, 0.7, 0.08, 0.5)
    _tb(slide, title, 1.0, 0.6, 11.33, 0.7,
        size=28, bold=True, color=pal["text"], font=pal["title_font"])

    y = 1.6
    for bullet in bullets[:5]:
        # Accent marker
        _tb(slide, "▸", 1.0, y, 0.3, 0.5,
            size=18, color=pal["accent"], bold=True, font=pal["body_font"])
        _tb(slide, bullet, 1.4, y, 10.5, 0.6,
            size=16, color=pal["text"], font=pal["body_font"])
        y += 0.65

    if callout:
        bar_top = 5.5
        _accent_bar(slide, pal, 0.8, bar_top, 0.08, 0.8)
        _tb(slide, "KEY TAKEAWAY", 1.0, bar_top, 11.33, 0.3,
            size=11, color=pal["accent"], bold=True, font=pal["body_font"])
        _tb(slide, callout, 1.0, bar_top + 0.3, 11.33, 0.7,
            size=15, color=pal["text"], font=pal["body_font"])

    if tools:
        y = 6.6
        x = 1.0
        for tool in tools[:6]:
            chip = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(len(tool) * 0.09 + 0.3), Inches(0.35),
            )
            chip.fill.solid(); chip.fill.fore_color.rgb = pal["surface"]
            chip.line.color.rgb = pal["subtext"]
            chip_tf = chip.text_frame
            chip_tf.margin_left = chip_tf.margin_right = Inches(0.1)
            chip_tf.margin_top = chip_tf.margin_bottom = Inches(0.02)
            p = chip_tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = tool
            run.font.size = Pt(11)
            run.font.color.rgb = pal["subtext"]
            run.font.name = pal["body_font"]
            x += len(tool) * 0.09 + 0.5
    return slide


def comparison_slide(prs, title: str, left_header: str, left_items: list[str],
                     right_header: str, right_items: list[str], style: str = "syndara"):
    """Two-column comparison: e.g. Before | After, Manual | AI-assisted."""
    pal = _palette(style)
    slide = _blank_slide(prs, pal)
    _tb(slide, title, 1.0, 0.6, 11.33, 0.7,
        size=28, bold=True, color=pal["text"], align=PP_ALIGN.CENTER, font=pal["title_font"])

    col_w = 5.2
    # Left column
    _tb(slide, left_header, 1.0, 1.8, col_w, 0.5,
        size=18, bold=True, color=pal["accent"], align=PP_ALIGN.CENTER, font=pal["title_font"])
    for i, item in enumerate(left_items[:4]):
        _tb(slide, f"• {item}", 1.2, 2.5 + i * 0.75, col_w - 0.4, 0.65,
            size=14, color=pal["text"], font=pal["body_font"])

    # Right column
    _tb(slide, right_header, 7.13, 1.8, col_w, 0.5,
        size=18, bold=True, color=pal["accent"], align=PP_ALIGN.CENTER, font=pal["title_font"])
    for i, item in enumerate(right_items[:4]):
        _tb(slide, f"• {item}", 7.33, 2.5 + i * 0.75, col_w - 0.4, 0.65,
            size=14, color=pal["text"], font=pal["body_font"])

    # Vertical divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.63), Inches(1.8), Inches(0.04), Inches(4.5),
    )
    div.fill.solid(); div.fill.fore_color.rgb = pal["subtext"]; div.line.fill.background()
    return slide


def stats_slide(prs, title: str, stats: list[dict], style: str = "syndara"):
    """Big-number stats highlight. stats is list of {number, label}."""
    pal = _palette(style)
    slide = _blank_slide(prs, pal)
    _tb(slide, title, 1.0, 0.6, 11.33, 0.7,
        size=28, bold=True, color=pal["text"], align=PP_ALIGN.CENTER, font=pal["title_font"])

    n = min(len(stats), 3)
    if n == 0:
        return slide
    col_w = 11.33 / n
    for i, stat in enumerate(stats[:n]):
        x = col_w * i + 0.5
        _tb(slide, str(stat.get("number", "")), x, 2.5, col_w - 1.0, 1.8,
            size=80, bold=True, color=pal["accent"], align=PP_ALIGN.CENTER, font=pal["title_font"])
        _tb(slide, str(stat.get("label", "")), x, 4.5, col_w - 1.0, 1.0,
            size=14, color=pal["subtext"], align=PP_ALIGN.CENTER, font=pal["body_font"])
    return slide


def quote_slide(prs, quote: str, attribution: str = "", style: str = "syndara"):
    """Pullquote — large quote + attribution."""
    pal = _palette(style)
    slide = _blank_slide(prs, pal)
    _tb(slide, '"', 0.8, 1.5, 1.5, 2.0,
        size=120, bold=True, color=pal["accent"], font=pal["title_font"])
    _tb(slide, quote, 2.2, 2.5, 10.0, 2.5,
        size=26, color=pal["text"], font=pal["title_font"])
    if attribution:
        _tb(slide, "— " + attribution, 2.2, 5.3, 10.0, 0.6,
            size=16, color=pal["subtext"], font=pal["body_font"])
    return slide


def steps_slide(prs, title: str, steps: list[str], style: str = "syndara"):
    """Numbered step-by-step process."""
    pal = _palette(style)
    slide = _blank_slide(prs, pal)
    _accent_bar(slide, pal, 0.8, 0.7, 0.08, 0.5)
    _tb(slide, title, 1.0, 0.6, 11.33, 0.7,
        size=28, bold=True, color=pal["text"], font=pal["title_font"])

    n = min(len(steps), 5)
    y = 1.8
    for i, step in enumerate(steps[:n]):
        # Circled number
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.0), Inches(y), Inches(0.6), Inches(0.6),
        )
        circle.fill.solid(); circle.fill.fore_color.rgb = pal["accent"]
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.margin_left = ctf.margin_right = Inches(0)
        ctf.margin_top = ctf.margin_bottom = Inches(0)
        p = ctf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = str(i + 1)
        run.font.size = Pt(20); run.font.bold = True
        run.font.color.rgb = pal["bg"]; run.font.name = pal["title_font"]
        # Step text
        _tb(slide, step, 1.8, y + 0.05, 10.0, 0.6,
            size=16, color=pal["text"], font=pal["body_font"])
        y += 0.95
    return slide


def summary_slide(prs, title: str, takeaways: list[str], cta: str = "", style: str = "syndara"):
    """Closing slide — key takeaways + optional call to action."""
    pal = _palette(style)
    slide = _blank_slide(prs, pal)
    _accent_bar(slide, pal, 0.8, 0.7, 0.08, 0.5)
    _tb(slide, title, 1.0, 0.6, 11.33, 0.7,
        size=28, bold=True, color=pal["text"], font=pal["title_font"])

    y = 1.8
    for takeaway in takeaways[:4]:
        _tb(slide, "✓", 1.0, y, 0.4, 0.55,
            size=20, color=pal["accent"], bold=True, font=pal["title_font"])
        _tb(slide, takeaway, 1.5, y, 10.5, 0.7,
            size=16, color=pal["text"], font=pal["body_font"])
        y += 0.75

    if cta:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3.0), Inches(5.8), Inches(7.33), Inches(0.9),
        )
        bar.fill.solid(); bar.fill.fore_color.rgb = pal["accent"]; bar.line.fill.background()
        tf = bar.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = cta
        run.font.size = Pt(18); run.font.bold = True
        run.font.color.rgb = pal["bg"]; run.font.name = pal["title_font"]
    return slide


# Registry for the agent to call by name via a single tool.
LAYOUTS = {
    "title_slide": title_slide,
    "section_divider": section_divider,
    "agenda_slide": agenda_slide,
    "bullet_slide": bullet_slide,
    "comparison_slide": comparison_slide,
    "stats_slide": stats_slide,
    "quote_slide": quote_slide,
    "steps_slide": steps_slide,
    "summary_slide": summary_slide,
}


def add_slide(prs, layout_name: str, speaker_notes: str = "", style: str = "professional", **kwargs):
    """Dispatch to the named layout function. Main entry point for the agent."""
    if layout_name not in LAYOUTS:
        raise ValueError(
            f"Unknown layout '{layout_name}'. Available: {', '.join(LAYOUTS.keys())}"
        )
    slide = LAYOUTS[layout_name](prs, style=style, **kwargs)
    if speaker_notes:
        _set_notes(slide, speaker_notes)
    return slide
