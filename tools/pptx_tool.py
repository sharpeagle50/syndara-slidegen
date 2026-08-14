"""PPTX tool: creates and modifies PowerPoint files via python-pptx."""
from __future__ import annotations
import json
import re
import subprocess
import sys
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Syndara brand colors
COLOR_BG = RGBColor(0x0F, 0x17, 0x2A)       # dark navy
COLOR_ACCENT = RGBColor(0x4A, 0xD8, 0xC4)   # teal
COLOR_TEXT = RGBColor(0xFF, 0xFF, 0xFF)      # white
COLOR_SUBTEXT = RGBColor(0xB0, 0xBE, 0xC5)  # light gray


def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size=18, bold=False, color=COLOR_TEXT,
                  align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def create_pptx(outline: dict, output_path: str) -> str:
    """
    Create a .pptx file from a slide outline dict.
    outline shape:
    {
        "title": "Course Title",
        "module_id": "...",
        "slides": [
            {
                "title": "Slide Title",
                "key_points": ["point 1", "point 2"],
                "speaker_notes": "Full narration text...",
                "tools_mentioned": ["ChatGPT", "Claude"],
                "callout": "Optional highlighted callout text"
            }
        ],
        "citations": ["Author (Year). Title. URL"]
    }
    Returns the output_path on success.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    # --- Title slide ---
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, COLOR_BG)
    _add_text_box(
        slide, outline.get("title", "Untitled Course"),
        Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
        font_size=40, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER
    )
    module_id = outline.get("module_id", "")
    if module_id:
        _add_text_box(
            slide, f"Module: {module_id}",
            Inches(1), Inches(4.2), Inches(11.33), Inches(0.5),
            font_size=18, color=COLOR_SUBTEXT, align=PP_ALIGN.CENTER
        )
    _add_text_box(
        slide, "Built with Syndara · Ava Zarkesh · Funded by AyZar Outreach",
        Inches(1), Inches(6.5), Inches(11.33), Inches(0.5),
        font_size=12, color=COLOR_SUBTEXT, align=PP_ALIGN.CENTER
    )
    # No speaker notes on title slide

    # --- Content slides ---
    for i, s in enumerate(outline.get("slides", []), 1):
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, COLOR_BG)

        # Slide number + title
        _add_text_box(
            slide, f"{i}. {s.get('title', '')}",
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.9),
            font_size=28, bold=True, color=COLOR_ACCENT
        )

        # Tools mentioned badge
        tools = s.get("tools_mentioned", [])
        if tools:
            _add_text_box(
                slide, "Tools: " + " · ".join(tools),
                Inches(0.5), Inches(1.1), Inches(12), Inches(0.35),
                font_size=12, color=COLOR_SUBTEXT
            )

        # Key points
        key_points = s.get("key_points", [])
        y_start = Inches(1.55)
        for j, point in enumerate(key_points[:5]):
            _add_text_box(
                slide, f"• {point}",
                Inches(0.6), y_start + Inches(j * 0.75), Inches(11.5), Inches(0.7),
                font_size=18, color=COLOR_TEXT
            )

        # Callout box
        callout = s.get("callout", "")
        if callout:
            _add_text_box(
                slide, f"▶  {callout}",
                Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.0),
                font_size=16, bold=True, color=COLOR_BG,
            )
            # Add colored bg rectangle behind callout
            from pptx.util import Emu
            txBox = slide.shapes[-1]
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_ACCENT
            shape.line.fill.background()
            # Move text box in front of rectangle
            slide.shapes._spTree.remove(txBox._element)
            slide.shapes._spTree.append(txBox._element)

        # Speaker notes
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = s.get("speaker_notes", "")

    # --- Citations slide ---
    citations = outline.get("citations", [])
    if citations:
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, COLOR_BG)
        _add_text_box(
            slide, "Sources & Citations",
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
            font_size=28, bold=True, color=COLOR_ACCENT
        )
        for j, cite in enumerate(citations[:10]):
            _add_text_box(
                slide, f"[{j+1}] {cite}",
                Inches(0.6), Inches(1.2 + j * 0.55), Inches(12), Inches(0.5),
                font_size=13, color=COLOR_SUBTEXT
            )

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def insert_image(pptx_path: str, slide_idx: int, image_path: str,
                 left=None, top=None, width=None, height=None) -> str:
    """Insert a PNG image into a specific slide."""
    from pptx.util import Inches
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]
    left = left or Inches(7.5)
    top = top or Inches(1.5)
    width = width or Inches(5.0)
    slide.shapes.add_picture(image_path, left, top, width=width)
    prs.save(pptx_path)
    return pptx_path


def update_speaker_notes(pptx_path: str, slide_idx: int, notes: str) -> str:
    """Update speaker notes on a specific slide."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    prs.save(pptx_path)
    return pptx_path


# ── Progressive reveal ("build") support ──────────────────────────────────────
# A builder marks a slide for a progressive build by (1) tagging animated shapes with an objectName
# of "reveal:<ranges>" where <ranges> is comma-separated visibility windows over beat numbers:
#   reveal:2       visible from beat 2 to the end
#   reveal:1-1     visible only during beat 1 (a transient overlay)
#   reveal:1-1,3   in at beat 1, out, back from beat 3 onward (re-entry)
# Untagged shapes are always visible. And (2) writing the speaker notes as a NARRATION SCRIPT with
# inline stage-direction markers `[[ N | what happens ]]` after a leading "[REVEAL]" header — each
# marker sits at the exact word beat N starts, and cutting the script at the markers yields one
# narration segment per beat (segment 0 = the base, before the first marker).
#
# explode_build_slides() then turns each such slide into V = M+1 physical slides IN PLACE: the
# original slide is duplicated per beat and shapes not visible at that beat get the native OOXML
# hidden="1" attribute (PowerPoint's Selection-Pane hide — respected by PowerPoint, LibreOffice's
# renderer, and Google Slides). Nothing is deleted and z-order is untouched, so every beat is
# pixel-identical to the authored slide except which elements show, and collapse_build_slides() is
# a lossless inverse (drop the earlier beats, unhide, rejoin the notes script). Each beat's notes
# carry its own clean narration segment plus one ASCII cue line `(build G | beat B/V | desc)` that
# doubles as a presenter stage cue and the round-trip metadata.
_REVEAL_NAME_RE = re.compile(r"^reveal:(\d+(?:-\d+)?(?:,\d+(?:-\d+)?)*)$", re.IGNORECASE)
_REVEAL_MARKER_RE = re.compile(r"\[\[[^\[\]]*\]\]")   # any inline stage direction [[ … ]]
_REVEAL_MARKER_STEP_RE = re.compile(r"\[\[\s*(\d+)\s*\|\s*([^\[\]]*?)\s*\]\]")  # [[ N | desc ]]
_BUILD_CUE_RE = re.compile(r"^\(build (\d+) \| beat (\d+)/(\d+)(?: \| ([^)\n]*))?\)\s*$", re.MULTILINE)
MAX_BUILD_BEATS = 9   # sanity cap: at most 9 markers (10 physical slides) per build slide


def _parse_reveal_segments(notes: str) -> list[str]:
    """Split reveal speaker notes into ordered per-step narration segments by cutting at the inline
    [[ … ]] markers. Segment 0 is the text before the first marker (base/step-0 narration); each
    following segment is what's narrated while its step's element(s) are on screen. Returns [] when
    the notes are not a reveal block (no leading [REVEAL]) so normal slides are unaffected."""
    if not notes or "[REVEAL]" not in notes:
        return []
    body = notes.split("[REVEAL]", 1)[1]
    parts = [p.strip() for p in _REVEAL_MARKER_RE.split(body)]
    return [p for p in parts if p]


def strip_reveal_markup(notes: str) -> str:
    """Clean narration with the [REVEAL] header and every inline [[ … ]] stage-direction marker
    removed — what the downloadable pptx speaker notes should show (no animation directions)."""
    if not notes:
        return notes
    text = _REVEAL_MARKER_RE.sub(" ", notes.replace("[REVEAL]", " "))
    return re.sub(r"\s+", " ", text).strip()


def clean_reveal_notes(pptx_path: str) -> int:
    """Rewrite every reveal slide's speaker notes to clean narration (markers stripped), so the
    downloadable pptx reads normally. Returns the number of slides cleaned; no-op on slides without
    reveal markup. Call AFTER the reveal frames/segments have been extracted from this pptx."""
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return 0
    n = 0
    for slide in prs.slides:
        if not slide.has_notes_slide:
            continue
        tf = slide.notes_slide.notes_text_frame
        # Not gated on the [REVEAL] header: headerless "[[N | ...]]" markers (a model slip the
        # strict script parser rejects) must not survive into the downloaded deck's notes either.
        if tf.text and ("[REVEAL]" in tf.text or _REVEAL_MARKER_RE.search(tf.text)):
            tf.text = strip_reveal_markup(tf.text)
            n += 1
    if n:
        prs.save(pptx_path)
    return n


def _reveal_ranges(name: str) -> Optional[list[tuple[int, Optional[int]]]]:
    """Parse a reveal objectName into visibility windows [(first, last|None), ...] — None if the
    shape isn't reveal-tagged. A bare number means "from that beat to the end"."""
    m = _REVEAL_NAME_RE.match(name or "")
    if not m:
        return None
    out: list[tuple[int, Optional[int]]] = []
    for part in m.group(1).split(","):
        if "-" in part:
            a, b = part.split("-")
            out.append((int(a), int(b)))
        else:
            out.append((int(part), None))
    return out


def _shape_visible_at(ranges: list[tuple[int, Optional[int]]], k: int) -> bool:
    """A reveal-tagged shape is visible on beat k iff k falls inside any of its windows."""
    return any(first <= k and (last is None or k <= last) for first, last in ranges)


def parse_reveal_script(notes: str) -> Optional[dict]:
    """Parse a [REVEAL] narration script into {"segments": [seg0..segM], "markers": [(1,desc1)..]}.
    Returns None unless well-formed: a [REVEAL] header, 1..MAX_BUILD_BEATS markers whose step
    numbers are exactly 1,2,...,M in order, and a non-empty narration segment for every beat
    (that's what gives each beat its timing/narration). Malformed → None → static slide."""
    if not notes or "[REVEAL]" not in notes:
        return None
    body = notes.split("[REVEAL]", 1)[1]
    markers = [(int(m.group(1)), m.group(2).strip()) for m in _REVEAL_MARKER_STEP_RE.finditer(body)]
    # Any [[…]] not matching the N|desc form makes the script ambiguous — reject.
    if len(markers) != len(_REVEAL_MARKER_RE.findall(body)):
        return None
    M = len(markers)
    if not (1 <= M <= MAX_BUILD_BEATS) or [k for k, _ in markers] != list(range(1, M + 1)):
        return None
    segments = [p.strip() for p in _REVEAL_MARKER_RE.split(body)]
    if len(segments) != M + 1 or any(not s for s in segments):
        return None
    return {"segments": segments, "markers": markers}


def _slide_notes(slide) -> str:
    return slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""


def _set_slide_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def _set_shape_hidden(shp, hidden: bool) -> None:
    """Toggle the native OOXML hidden attribute (p:cNvPr@hidden) — for every shape kind
    (sp/pic/graphicFrame/grpSp/cxnSp) the cNvPr is the first child of the first child."""
    cNvPr = shp._element[0][0]
    if hidden:
        cNvPr.set("hidden", "1")
    else:
        cNvPr.attrib.pop("hidden", None)


def _duplicate_slide(prs, index: int):
    """Insert an exact duplicate of prs.slides[index] immediately BEFORE it and return it.
    Shape XML is deep-copied; relationships (images/charts/media) are re-created on the new part
    pointing at the SAME underlying parts (no file-size blowup), with r:id references rewritten
    when the new part assigns a different rId."""
    import copy as _copy
    from pptx.oxml.ns import qn
    src = prs.slides[index]
    dup = prs.slides.add_slide(src.slide_layout)
    # Drop the placeholder shapes the layout seeded, then copy every shape element verbatim.
    for shp in list(dup.shapes):
        shp._element.getparent().remove(shp._element)
    for el in list(src._element.spTree)[2:]:      # [0]=nvGrpSpPr [1]=grpSpPr, rest are shapes
        dup._element.spTree.append(_copy.deepcopy(el))
    # Slide-level background, if any.
    src_bg = src._element.cSld.find(qn("p:bg"))
    if src_bg is not None and dup._element.cSld.find(qn("p:bg")) is None:
        dup._element.cSld.insert(0, _copy.deepcopy(src_bg))
    # Re-create the source slide's relationships on the duplicate.
    id_map: dict[str, str] = {}
    for rId, rel in src.part.rels.items():
        # Skip the layout (the duplicate already has its own) and the NOTES relationship — copying
        # it would make every beat share the source's one notes part, so the last notes write would
        # win on all of them. Leaving it out lets each beat lazily create its own notes part.
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_rId = dup.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rId = dup.part.relate_to(rel.target_part, rel.reltype)
        if new_rId != rId:
            id_map[rId] = new_rId
    if id_map:
        RNS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
        for el in dup._element.iter():
            for attr, val in list(el.attrib.items()):
                if attr.startswith(RNS) and val in id_map:
                    el.attrib[attr] = id_map[val]
    # add_slide appended at the end — move the reference to sit right before the source slide.
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    new_id = ids[-1]
    sldIdLst.remove(new_id)
    sldIdLst.insert(index, new_id)
    return dup


def _build_cue(group: int, beat: int, total: int, desc: str = "") -> str:
    return f"(build {group} | beat {beat}/{total}" + (f" | {desc}" if desc else "") + ")"


def explode_build_slides(pptx_path: str) -> int:
    """Expand every well-formed build slide into its beats, IN PLACE. Beat k is a duplicate of the
    authored slide with shapes not visible at k carrying hidden="1"; its notes are that beat's clean
    narration segment plus a `(build G | beat B/V | desc)` cue line. The authored slide itself
    becomes the last beat (its transients hidden too). Returns the number of physical slides added.
    Slides with no reveal tags or a malformed script are left untouched (static)."""
    prs = Presentation(pptx_path)
    targets = []   # (index, script, tagged_ranges_ok)
    for i, slide in enumerate(prs.slides):
        script = parse_reveal_script(_slide_notes(slide))
        if not script:
            continue
        M = len(script["markers"])
        tag_lists = [r for r in (_reveal_ranges(sh.name) for sh in slide.shapes) if r]
        if not tag_lists:
            continue
        if any(first > M or (last is not None and last > M) for rl in tag_lists for first, last in rl):
            continue   # a tag references a beat past the last marker — malformed, stay static
        targets.append((i, script))
    if not targets:
        return 0
    group_of = {i: g + 1 for g, (i, _) in enumerate(targets)}
    added = 0
    for i, script in reversed(targets):   # reverse so insertions don't shift pending indices
        g = group_of[i]
        M = len(script["markers"])
        V = M + 1
        for _ in range(M):                # M duplicates before the original → beats 0..M-1
            _duplicate_slide(prs, i)
            added += 1
        for k in range(V):                # beats live at indices i..i+M; original is last (beat M)
            beat_slide = prs.slides[i + k]
            for shp in beat_slide.shapes:
                rl = _reveal_ranges(shp.name)
                if rl is not None:
                    _set_shape_hidden(shp, not _shape_visible_at(rl, k))
            desc = script["markers"][k - 1][1] if k >= 1 else ""
            _set_slide_notes(beat_slide, script["segments"][k] + "\n\n" + _build_cue(g, k + 1, V, desc))
    prs.save(pptx_path)
    return added


def collapse_build_slides(pptx_path: str) -> int:
    """Lossless inverse of explode_build_slides, for re-ingesting a downloaded/edited deck: drop
    every non-final beat, un-hide all shapes on the kept slide, and rebuild the original [REVEAL]
    narration script from the beats' notes + cue descriptions. Slides without build cues are
    untouched; incomplete groups (user deleted beats in PowerPoint) keep their last surviving beat.
    Returns the number of physical slides removed."""
    from pptx.oxml.ns import qn
    prs = Presentation(pptx_path)
    # Collect group membership in slide order.
    members: dict[int, list[tuple[int, int, int, str, str]]] = {}   # g -> [(idx, beat, total, desc, segment)]
    for i, slide in enumerate(prs.slides):
        notes = _slide_notes(slide)
        m = _BUILD_CUE_RE.search(notes)
        if not m:
            continue
        seg = _BUILD_CUE_RE.sub("", notes).strip()
        members.setdefault(int(m.group(1)), []).append(
            (i, int(m.group(2)), int(m.group(3)), (m.group(4) or "").strip(), seg))
    if not members:
        return 0
    to_remove: list[int] = []
    for g, beats in members.items():
        beats.sort(key=lambda b: b[1])
        final_idx = beats[-1][0]          # highest surviving beat = best-effort final
        final_slide = prs.slides[final_idx]
        for shp in final_slide.shapes:
            _set_shape_hidden(shp, False)
        # Rebuild the narration script: seg0 [[1|d1]] seg1 ... using each beat's segment + desc.
        script = "[REVEAL] " + beats[0][4]
        for _, beat_no, _, desc, seg in beats[1:]:
            script += f" [[{beat_no - 1} | {desc}]] {seg}"
        _set_slide_notes(final_slide, script if len(beats) > 1 else beats[0][4])
        to_remove.extend(idx for idx, *_ in beats[:-1])
    # Drop the removed slides' parts + references (reverse order keeps indices valid).
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    for idx in sorted(to_remove, reverse=True):
        sld_id = ids[idx]
        rId = sld_id.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sld_id)
    prs.save(pptx_path)
    return len(to_remove)


# DrawingML preset-geometry (ST_ShapeType) aliases an LLM commonly emits that
# are NOT valid enum values. An invalid prst= value is a schema violation, not
# a structural one: LibreOffice silently drops the shape (invisible), while
# PowerPoint throws the "repair" dialog and replaces it with a degenerate
# placeholder (a diagonal line). The canonical example is `oval` — the real
# token is `ellipse`. Normalize the common offenders to their valid tokens.
# Keys are scoped to `prstGeom prst="..."` (not a bare `prst="..."`) so the
# replace can only ever touch a shape's geometry attribute — never a text run,
# a preset text-warp, or any other element that happens to contain the string.
_PRESET_ALIASES = {
    b'prstGeom prst="oval"': b'prstGeom prst="ellipse"',
    b'prstGeom prst="circle"': b'prstGeom prst="ellipse"',
    b'prstGeom prst="rectangle"': b'prstGeom prst="rect"',
    b'prstGeom prst="square"': b'prstGeom prst="rect"',
    b'prstGeom prst="roundedRect"': b'prstGeom prst="roundRect"',
    b'prstGeom prst="roundRectangle"': b'prstGeom prst="roundRect"',
    b'prstGeom prst="roundedRectangle"': b'prstGeom prst="roundRect"',
}


def _fix_preset_aliases(data: bytes) -> tuple[bytes, bool]:
    """Replace invalid DrawingML preset-geometry names with their valid tokens."""
    changed = False
    for bad, good in _PRESET_ALIASES.items():
        if bad in data:
            data = data.replace(bad, good)
            changed = True
    return data, changed


def _fix_chart_value_axis_titles(data: bytes) -> tuple[bytes, bool]:
    """Rotate native-chart value-axis (y) titles to vertical.

    PptxGenJS emits value-axis titles with no rotation (<a:bodyPr/>). PowerPoint
    then renders the y-axis title HORIZONTAL and detached from the axis (its
    default for a missing rot); LibreOffice happens to default it vertical, which
    hides the bug in QA renders. Make the rotation explicit — rot="-5400000"
    (-90deg, the exact value PowerPoint itself writes for a y-axis title) — on
    each value-axis title's bodyPr. Category-axis (x) titles are deliberately
    left horizontal, so the patch is scoped to <c:valAx> blocks only.
    """
    changed = False

    def _patch_valax(block_m: "re.Match[bytes]") -> bytes:
        nonlocal changed

        def _add_rot(title_m: "re.Match[bytes]") -> bytes:
            nonlocal changed
            if b"rot=" in title_m.group(3):       # already rotated — leave it
                return title_m.group(0)
            changed = True
            return (title_m.group(1) + title_m.group(2) + title_m.group(3)
                    + b' rot="-5400000"' + title_m.group(4))

        # only the first <a:bodyPr> inside THIS value axis's <c:title>
        return re.sub(rb"(<c:title>.*?)(<a:bodyPr)([^>]*?)(/?>)",
                      _add_rot, block_m.group(0), count=1, flags=re.S)

    new = re.sub(rb"<c:valAx>.*?</c:valAx>", _patch_valax, data, flags=re.S)
    return new, changed


def _fix_chart_orphan_axids(data: bytes) -> tuple[bytes, bool]:
    """Drop chart axis references that point to no defined axis.

    PptxGenJS emits native 2-D bar charts whose <c:barChart> lists THREE
    <c:axId> entries while only two axes (<c:catAx> + <c:valAx>) are defined —
    the third id matches nothing. PowerPoint requires every axId in a chart
    group to resolve to a defined axis, so it flags the deck as damaged and
    shows the 'repair' dialog (and drops the chart). python-pptx, LibreOffice,
    and our repair-scanner all tolerate the orphan, which is why it survives QA
    and the structural checks. Remove any <c:axId val="N"/> whose N is not the
    id of a defined axis (a real serAx/dateAx keeps its id, so 3-D charts and
    secondary axes are untouched).
    """
    defined = set(re.findall(
        rb"<c:(?:catAx|valAx|dateAx|serAx)>\s*<c:axId val=\"(\d+)\"", data))
    if not defined:
        return data, False
    referenced = set(re.findall(rb"<c:axId val=\"(\d+)\"\s*/>", data))
    orphans = referenced - defined
    if not orphans:
        return data, False
    new = data
    for oid in orphans:
        new = re.sub(rb"<c:axId val=\"" + re.escape(oid) + rb"\"\s*/>", b"", new)
    return new, new != data


def _fix_degenerate_shape_xfrms(data: bytes) -> tuple[bytes, bool]:
    """Give zero/negative-area shapes a valid positive extent.

    PptxGenJS draws straight horizontal/vertical lines as <p:sp> autoshapes
    whose <a:ext> has a zero dimension (cy=0 for a horizontal line, cx=0 for a
    vertical one) and occasionally a negative one. PowerPoint rejects a
    zero/negative-area shape extent and shows the 'repair' dialog; python-pptx,
    LibreOffice, and our scanner all tolerate it (LibreOffice silently rewrites
    it on re-save, which is why a round-tripped copy opens cleanly). Normalize
    each SHAPE xfrm: flip negative extents (preserving geometry via off +
    flipH/flipV) and clamp any remaining zero dimension to 1 EMU. The root
    group's xfrm (ext 0x0, followed by chOff/chExt) is matched-out because its
    <a:ext> is not immediately followed by </a:xfrm>, so the legitimate 0x0
    group bounds are left untouched.
    """
    changed = False

    def _norm(m: "re.Match[bytes]") -> bytes:
        nonlocal changed
        attrs = m.group(1).decode()
        x, y, cx, cy = (int(m.group(i)) for i in (2, 3, 4, 5))
        fh = 'flipH="1"' in attrs
        fv = 'flipV="1"' in attrs
        rot = re.search(r'rot="(-?\d+)"', attrs)
        nx, ny, ncx, ncy, nfh, nfv = x, y, cx, cy, fh, fv
        if ncx < 0:
            nx += ncx; ncx = -ncx; nfh = not nfh
        if ncy < 0:
            ny += ncy; ncy = -ncy; nfv = not nfv
        if ncx == 0:
            ncx = 1
        if ncy == 0:
            ncy = 1
        if (nx, ny, ncx, ncy, nfh, nfv) == (x, y, cx, cy, fh, fv):
            return m.group(0)
        changed = True
        a = ""
        if rot:
            a += f' rot="{rot.group(1)}"'
        if nfh:
            a += ' flipH="1"'
        if nfv:
            a += ' flipV="1"'
        return (f'<a:xfrm{a}><a:off x="{nx}" y="{ny}"/>'
                f'<a:ext cx="{ncx}" cy="{ncy}"/></a:xfrm>').encode()

    new = re.sub(
        rb'<a:xfrm([^>]*)><a:off x="(-?\d+)" y="(-?\d+)"/>'
        rb'<a:ext cx="(-?\d+)" cy="(-?\d+)"/></a:xfrm>',
        _norm, data)
    return new, changed


def _split_shared_master_themes(raw: bytes) -> tuple[bytes, bool]:
    """Give every master its own theme part.

    PowerPoint requires each master (slide / notes / handout) to reference a
    DISTINCT theme part. PptxGenJS emits a single theme1.xml shared by the slide
    master and the notes master; PowerPoint rejects the shared reference and
    shows the 'repair' dialog (then strips the notes master), while python-pptx,
    LibreOffice, and our scanner all tolerate it — LibreOffice silently clones
    the theme on re-save, which is why a round-tripped copy opens cleanly. When a
    theme is referenced by more than one master, clone it so each extra master
    gets its own copy, repoint that master's relationship, and register the new
    part's content type. Idempotent: a deck whose masters already have distinct
    themes is returned unchanged.
    """
    import io
    import posixpath
    import zipfile

    src = zipfile.ZipFile(io.BytesIO(raw))
    names = src.namelist()
    master_rels = [n for n in names if re.match(
        r"ppt/(?:slide|notes|handout)Masters/_rels/[^/]+\.xml\.rels$", n)]
    theme_users: "dict[str, list[tuple[str, str]]]" = {}
    for rn in master_rels:
        text = src.read(rn).decode("utf-8", "replace")
        base = posixpath.dirname(posixpath.dirname(rn))
        for m in re.finditer(r"<Relationship\b([^>]*)/?>", text):
            attrs = dict(re.findall(r'(\w+)="([^"]*)"', m.group(1)))
            if attrs.get("Type", "").endswith("/theme"):
                tgt = posixpath.normpath(posixpath.join(base, attrs.get("Target", "")))
                theme_users.setdefault(tgt, []).append((rn, attrs.get("Id", "")))

    shared = {t: u for t, u in theme_users.items() if len(u) > 1}
    if not shared:
        return raw, False

    used = [int(mm.group(1)) for n in names
            for mm in [re.match(r"ppt/theme/theme(\d+)\.xml$", n)] if mm]
    idx = (max(used) + 1) if used else 1

    new_parts: "dict[str, bytes]" = {}
    rels_edits: "dict[str, list[tuple[str, str]]]" = {}
    ct_adds: "list[str]" = []
    for theme_part, users in shared.items():
        body = src.read(theme_part)
        for rels_name, rel_id in users[1:]:   # the first master keeps the original
            new_theme = f"ppt/theme/theme{idx}.xml"
            idx += 1
            new_parts[new_theme] = body
            ct_adds.append(new_theme)
            base = posixpath.dirname(posixpath.dirname(rels_name))
            rels_edits.setdefault(rels_name, []).append(
                (rel_id, posixpath.relpath(new_theme, base)))

    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst:
        for info in src.infolist():
            data = src.read(info.filename)
            if info.filename in rels_edits:
                t = data.decode("utf-8", "replace")
                for rel_id, new_target in rels_edits[info.filename]:
                    t = re.sub(
                        r'(<Relationship\b[^>]*\bId="' + re.escape(rel_id)
                        + r'"[^>]*\bTarget=")[^"]*(")',
                        lambda mm, nt=new_target: mm.group(1) + nt + mm.group(2),
                        t, count=1)
                data = t.encode("utf-8")
            elif info.filename == "[Content_Types].xml":
                t = data.decode("utf-8", "replace")
                adds = "".join(
                    f'<Override PartName="/{p}" ContentType='
                    f'"application/vnd.openxmlformats-officedocument.theme+xml"/>'
                    for p in ct_adds)
                data = t.replace("</Types>", adds + "</Types>").encode("utf-8")
            zi = zipfile.ZipInfo(info.filename, date_time=info.date_time)
            zi.compress_type = info.compress_type
            dst.writestr(zi, data)
        for path, body in new_parts.items():
            dst.writestr(path, body)
    return out.getvalue(), True


def _strip_em_dashes_in_runs(data: bytes) -> tuple[bytes, bool]:
    """Replace Unicode dashes inside <a:t> text runs (visible slide text + speaker
    notes) with dash-free punctuation — the deterministic backstop for the few
    em dashes the model emits despite the style rule. Only edits text inside
    <a:t>, and only Unicode dashes, so markup/attributes/ASCII are never touched.
    """
    if not any(d in data for d in (b"\xe2\x80\x94", b"\xe2\x80\x93", b"\xe2\x80\x95")):
        return data, False
    try:
        from ..agents.base import strip_em_dashes
        text = data.decode("utf-8", "replace")
        new = re.sub(r"<a:t>(.*?)</a:t>",
                     lambda m: "<a:t>" + strip_em_dashes(m.group(1)) + "</a:t>",
                     text, flags=re.S)
        nb = new.encode("utf-8")
        return nb, nb != data
    except Exception:
        # Cosmetic only — must never disable the structural repairs around it.
        return data, False


_LEADING_PUNCT_RE = re.compile(r"^[,;:·]+\s+")


def _strip_stray_leading_punct(data: bytes) -> tuple[bytes, bool]:
    """Strip stray leading punctuation (", Linear IVP with a variable coefficient") from the
    FIRST text run of each paragraph — the deterministic backstop for a template artifact the
    model occasionally emits in titles. Visual QA classes it as critical placeholder_artifacts,
    so without this backstop a one-character defect buys a full builder rebuild pass.
    Paragraph-leading only: punctuation at the start of a LATER run in the same paragraph is
    legitimate mid-sentence text (runs are fragments) and is never touched.
    """
    if not any(p in data for p in (b"<a:t>,", b"<a:t>;", b"<a:t>:", b"<a:t>\xc2\xb7")):
        return data, False
    try:
        text = data.decode("utf-8", "replace")

        def _fix_para(pm: "re.Match[str]") -> str:
            return re.sub(
                r"<a:t>(.*?)</a:t>",
                lambda tm: "<a:t>" + _LEADING_PUNCT_RE.sub("", tm.group(1)) + "</a:t>",
                pm.group(0), count=1, flags=re.S,
            )

        new = re.sub(r"<a:p>.*?</a:p>", _fix_para, text, flags=re.S)
        nb = new.encode("utf-8")
        return nb, nb != data
    except Exception:
        # Cosmetic only — must never disable the structural repairs around it.
        return data, False


def _sanitize_package_bytes(raw: bytes, _depth: int = 0) -> tuple[bytes, bool]:
    """Rewrite an OOXML package (zip) to remove two classes of defect that make
    PowerPoint flag the file as damaged, recursing into embedded OOXML parts.
    Returns (bytes, changed).

    1. OPC-violating directory entries — PptxGenJS builds chart data workbooks
       with JSZip, which emits zero-byte 'folder' entries (names ending in '/').
       The OPC spec forbids directory parts, and they carry no content-type, so
       PowerPoint's strict parser flags the file as damaged and 'repairs' it
       (silently dropping the chart). python-pptx and LibreOffice are lenient
       and never surface this.
    2. Invalid preset-geometry names (e.g. prst="oval") — see _PRESET_ALIASES.
    """
    import io
    import zipfile

    src = zipfile.ZipFile(io.BytesIO(raw))
    changed = False
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst:
        for info in src.infolist():
            if info.filename.endswith("/"):
                changed = True  # drop directory entry
                continue
            data = src.read(info.filename)
            fname = info.filename.lower()
            if _depth < 3 and fname.endswith((".xlsx", ".docx", ".pptx")):
                data, sub_changed = _sanitize_package_bytes(data, _depth + 1)
                changed = changed or sub_changed
            elif fname.endswith(".xml"):
                data, preset_changed = _fix_preset_aliases(data)
                changed = changed or preset_changed
                if "slide" in fname:   # slides/layouts/masters/notesSlides
                    data, xfrm_changed = _fix_degenerate_shape_xfrms(data)
                    changed = changed or xfrm_changed
                    data, dash_changed = _strip_em_dashes_in_runs(data)
                    changed = changed or dash_changed
                    data, punct_changed = _strip_stray_leading_punct(data)
                    changed = changed or punct_changed
                if "charts/chart" in fname:
                    data, axis_changed = _fix_chart_value_axis_titles(data)
                    changed = changed or axis_changed
                    data, axid_changed = _fix_chart_orphan_axids(data)
                    changed = changed or axid_changed
            zi = zipfile.ZipInfo(info.filename, date_time=info.date_time)
            zi.compress_type = info.compress_type
            dst.writestr(zi, data)
    return (out.getvalue() if changed else raw), changed


def sanitize_pptx(pptx_path: str) -> bool:
    """Repair latent defects in a .pptx in place that make PowerPoint show the
    spurious 'PowerPoint found a problem with content / repaired and removed it'
    dialog (python-pptx and LibreOffice tolerate them all, so they never surface
    in our renders): OPC-violating directory entries (PptxGenJS chart workbooks),
    invalid preset geometry names (e.g. prst="oval"), chart value-axis title
    rotation, orphan chart axIds, zero/negative-area shape extents (straight
    lines), and a theme part shared by more than one master. Returns True if the
    file was modified. Safe and idempotent — a clean file is left untouched.
    """
    try:
        with open(pptx_path, "rb") as f:
            raw = f.read()
        cleaned, changed = _sanitize_package_bytes(raw)
        cleaned, theme_changed = _split_shared_master_themes(cleaned)
        changed = changed or theme_changed
        if changed:
            with open(pptx_path, "wb") as f:
                f.write(cleaned)
        return changed
    except Exception as e:
        print(f"[sanitize_pptx] skipped ({type(e).__name__}: {e})", flush=True)
        return False


def extract_slide_pngs(pptx_path: str, output_dir: str) -> list[str]:
    """
    Extract each slide as a PNG image.

    Uses PPTX → PDF (LibreOffice) → per-page PNGs (pdftoppm) for full-fidelity
    rendering of every slide. Falls back to matplotlib placeholders if
    LibreOffice or pdftoppm are unavailable.
    """
    import glob as globmod
    import shutil
    import tempfile

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = str(Path(pptx_path).resolve())

    # Step 1: PPTX → PDF via LibreOffice
    lo_candidates = [
        "soffice",
        "libreoffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/bin/libreoffice",
        "/usr/bin/soffice",
    ]
    # Nix store fallback — libreoffice-unwrapped buries the binary
    nix_hits = globmod.glob("/nix/store/*/bin/soffice")
    lo_candidates.extend(nix_hits)

    which_lo = shutil.which("soffice") or shutil.which("libreoffice")
    print(f"[extract_slide_pngs] which soffice/libreoffice → {which_lo}")
    print(f"[extract_slide_pngs] nix store hits: {nix_hits}")

    pdf_dir = tempfile.mkdtemp(prefix="syndara_pdf_")
    pdf_path = None
    try:
        # Placeholder-PNG policy: the matplotlib fallback exists for machines with NO conversion
        # tooling (local dev). When the tooling IS present and fails, falling back used to count as
        # SUCCESS all the way to publish — learners got "Preview unavailable" cards as their actual
        # slides, with no error anywhere. Now: tooling absent → fallback; tooling present but
        # failing → retry once, then RAISE so the module fails visibly and can be retried.
        lo_present = bool(which_lo or nix_hits
                          or any(Path(c).exists() for c in lo_candidates if c.startswith("/")))
        from .render_tool import libreoffice_convert_pdf
        _pdf = libreoffice_convert_pdf(pptx_path, lo_candidates, pdf_dir)
        if not _pdf and lo_present:
            print("[extract_slide_pngs] LibreOffice failed — retrying once", flush=True)
            _pdf = libreoffice_convert_pdf(pptx_path, lo_candidates, pdf_dir)
        if _pdf:
            pdf_path = Path(_pdf)
            print(f"[extract_slide_pngs] PDF created (isolated profile, scaled timeout)")

        if not pdf_path:
            if lo_present:
                raise RuntimeError(
                    f"LibreOffice is installed but failed to convert {pptx_path} after a retry — "
                    f"refusing to publish placeholder slide images")
            print("[extract_slide_pngs] LibreOffice not found — using matplotlib fallback")
            return _fallback_slide_pngs(pptx_path, str(output_dir))

        # Step 2: PDF → per-page PNGs via pdftoppm
        prefix = str(output_dir / "slide")
        for _attempt in (1, 2):
            try:
                result = subprocess.run(
                    # 200 DPI ≈ 2667px wide on a 13.33in slide — crisp on retina-class
                    # displays where the old 150 read soft next to PowerPoint's vector
                    # rendering (~1.8x the PNG bytes; R2 storage/bandwidth, acceptable).
                    ["pdftoppm", "-png", "-r", "200", str(pdf_path), prefix],
                    capture_output=True, timeout=180,
                )
                if result.returncode == 0:
                    break
                print(f"[extract_slide_pngs] pdftoppm failed (attempt {_attempt}): "
                      f"{result.stderr.decode()[:300]}")
                if _attempt == 2:
                    raise RuntimeError(
                        "pdftoppm failed twice on a valid PDF — refusing to publish placeholder "
                        "slide images")
            except FileNotFoundError:
                print("[extract_slide_pngs] pdftoppm not installed — using matplotlib fallback")
                return _fallback_slide_pngs(pptx_path, str(output_dir))
            except subprocess.TimeoutExpired:
                print(f"[extract_slide_pngs] pdftoppm timed out (attempt {_attempt})")
                if _attempt == 2:
                    raise RuntimeError(
                        "pdftoppm timed out twice — refusing to publish placeholder slide images")

    finally:
        shutil.rmtree(pdf_dir, ignore_errors=True)

    pngs = sorted(output_dir.glob("slide-*.png"))
    if not pngs:
        raise RuntimeError(
            f"Slide PNG conversion produced no images for {pptx_path} — refusing to publish "
            f"placeholder slide images")

    return [str(p) for p in pngs]


def _fallback_slide_pngs(pptx_path: str, output_dir: str) -> list[str]:
    """Last-resort placeholder PNGs when LibreOffice/pdftoppm unavailable."""
    from pptx import Presentation as Prs
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    out_path = Path(output_dir)
    prs = Prs(pptx_path)
    paths = []
    for i, slide in enumerate(prs.slides):
        fig, ax = plt.subplots(figsize=(13.33, 7.5), facecolor="#1E293B")
        ax.set_facecolor("#1E293B")
        ax.axis("off")

        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                title_text = shape.text.strip()[:80]
                break

        ax.text(
            0.5, 0.5, title_text or f"Slide {i+1}",
            transform=ax.transAxes, fontsize=24,
            color="#FFFFFF", ha="center", va="center",
            wrap=True, fontweight="bold",
        )
        ax.text(
            0.5, 0.35, "(Preview unavailable — install LibreOffice for full rendering)",
            transform=ax.transAxes, fontsize=10,
            color="#94A3B8", ha="center", va="center",
        )
        out = out_path / f"slide-{i+1:02d}.png"
        fig.savefig(str(out), dpi=150, bbox_inches="tight", facecolor="#1E293B")
        plt.close(fig)
        paths.append(str(out))

    return paths


def extract_deck_content(pptx_path: str, output_dir: str) -> list[dict]:
    """Extract text, speaker notes, and embedded images from every slide.

    Returns a list of dicts with slide_index, text, speaker_notes, and
    embedded_images (list of extracted image file paths with metadata).
    Tiny images (<100x100 px) are filtered as decorative.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    out_path = Path(output_dir)
    out_path.mkdir(parents=True, exist_ok=True)
    slides = []

    EMU_PER_INCH = 914400
    DPI = 96
    MIN_PX = 100

    def _iter_shapes(parent):
        """Yield all shapes, recursing into group shapes."""
        for shape in parent.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    yield from _iter_shapes(shape)
                except Exception:
                    pass
            else:
                yield shape

    for i, slide in enumerate(prs.slides):
        text_blocks = []
        embedded_images = []
        img_counter = 0

        for shape in _iter_shapes(slide):
            if shape.has_text_frame:
                text_blocks.append(shape.text_frame.text.strip())

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    w_px = int(shape.width / EMU_PER_INCH * DPI) if shape.width else 0
                    h_px = int(shape.height / EMU_PER_INCH * DPI) if shape.height else 0
                    if w_px < MIN_PX and h_px < MIN_PX:
                        continue
                    ext = image.ext or "png"
                    fname = f"slide_{i}_img_{img_counter}.{ext}"
                    img_path = out_path / fname
                    img_path.write_bytes(image.blob)
                    embedded_images.append({
                        "path": str(img_path),
                        "content_type": image.content_type or f"image/{ext}",
                        "width_px": w_px,
                        "height_px": h_px,
                    })
                    img_counter += 1
                except Exception as e:
                    print(f"[extract_deck_content] slide {i} image {img_counter}: {e}")
                    continue

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "slide_index": i,
            "text": "\n".join(t for t in text_blocks if t),
            "speaker_notes": notes,
            "embedded_images": embedded_images,
        })

    return slides


def extract_text_content(pptx_path: str) -> list[dict]:
    """Extract slide titles, bullets, and speaker notes as structured dicts."""
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        text_blocks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_blocks.append(shape.text_frame.text.strip())
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({
            "slide_index": i,
            "text": "\n".join(t for t in text_blocks if t),
            "speaker_notes": notes,
        })
    return slides
