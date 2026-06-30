"""Extract plain text from common document formats, for grounding deck content.

Used by the CLI's ``--context-file`` / ``--context-dir`` flags and shared with
Syndara's hosted upload handler so both sides have one implementation.

The heavy parsers (PDF / Word / Excel) are **optional dependencies** — install
them with ``pip install 'syndara-slidegen[context]'``. When one is missing, the
matching file type raises a clear error (or is skipped, for directory walks), so
the core package stays dependency-light. Deliberately framework-agnostic (no web
framework imports) so it can be reused anywhere.
"""
from __future__ import annotations

import io
from pathlib import Path

# Plain-text extensions read directly; binary formats need an optional parser.
# .ipynb is plain JSON, so it needs no extra dependency (handled inline below).
TEXT_EXTS = ("txt", "md", "csv", "tsv")
SUPPORTED_EXTS = set(TEXT_EXTS) | {"pdf", "docx", "pptx", "xlsx", "ipynb"}

_SUPPORTED_LABEL = "PDF, DOCX, PPTX, TXT, MD, CSV, Excel, Jupyter (.ipynb)"


class UnsupportedFileType(ValueError):
    """Raised when a file's extension isn't a supported text source."""


def extract_text(data: bytes, filename: str) -> str:
    """Return the text content of a document from its raw bytes + filename.

    Raises ``UnsupportedFileType`` for unknown extensions, and ``RuntimeError``
    with an install hint if the optional parser for a supported type is missing.
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in TEXT_EXTS:
        return data.decode("utf-8", errors="replace")
    if ext == "ipynb":
        return _extract_ipynb(data)
    if ext == "pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            raise RuntimeError("PDF support needs pypdf — install with: pip install 'syndara-slidegen[context]'")
        reader = PdfReader(io.BytesIO(data))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    if ext == "docx":
        try:
            from docx import Document
        except ImportError:
            raise RuntimeError("DOCX support needs python-docx — install with: pip install 'syndara-slidegen[context]'")
        doc = Document(io.BytesIO(data))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if ext == "pptx":
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError("PPTX support needs python-pptx — install with: pip install 'syndara-slidegen[context]'")
        return _extract_pptx(Presentation(io.BytesIO(data)))
    if ext == "xlsx":
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise RuntimeError("XLSX support needs openpyxl — install with: pip install 'syndara-slidegen[context]'")
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheets = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"[Sheet: {name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)
    raise UnsupportedFileType(f"Unsupported file type: .{ext}. Supported: {_SUPPORTED_LABEL}.")


def _extract_ipynb(data: bytes) -> str:
    """Flatten a Jupyter notebook (.ipynb, plain JSON) to text for grounding: markdown/raw
    cells as-is, code cells fenced. Cell outputs are skipped — they're often large or binary
    (images, tracebacks) and we want the authored content, not the run artifacts."""
    import json
    try:
        nb = json.loads(data.decode("utf-8", errors="replace"))
    except ValueError:
        raise UnsupportedFileType("Unsupported file type: .ipynb (not valid notebook JSON).")
    parts: list[str] = []
    for cell in nb.get("cells", []) or []:
        src = cell.get("source", "")
        text = ("".join(src) if isinstance(src, list) else str(src or "")).strip()
        if not text:
            continue
        parts.append(f"```python\n{text}\n```" if cell.get("cell_type") == "code" else text)
    return "\n\n".join(parts)


def _extract_pptx(prs) -> str:
    """Flatten a PowerPoint deck to text for grounding: per slide, the text from every
    shape (titles, bullets, text boxes), tables row-by-row (tab-separated), and any
    speaker notes. Shapes are emitted in document order; empty slides are skipped."""
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if text.strip():
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells]
                    if any(c.strip() for c in cells):
                        parts.append("\t".join(cells))
        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        if notes:
            parts.append(f"[Notes] {notes}")
        if parts:
            slides.append(f"[Slide {i}]\n" + "\n".join(parts))
    return "\n\n".join(slides)


def extract_text_from_path(path: str | Path) -> str:
    """Read a single file from disk and extract its text."""
    p = Path(path)
    return extract_text(p.read_bytes(), p.name)


def extract_text_from_dir(
    directory: str | Path,
    *,
    max_chars: int = 500_000,
    recursive: bool = False,
    log=lambda _m: None,
) -> str:
    """Concatenate extracted text from every supported file in a directory.

    Unsupported files (and supported files whose optional parser isn't
    installed) are skipped with a logged note rather than failing the whole run.
    Adds files in sorted order until ``max_chars`` is reached, so a huge folder
    can't blow up the context window.
    """
    d = Path(directory)
    if not d.is_dir():
        raise NotADirectoryError(f"Not a directory: {directory}")
    paths = sorted(d.rglob("*") if recursive else d.glob("*"))
    parts: list[str] = []
    total = 0
    for p in paths:
        if not p.is_file() or p.suffix.lstrip(".").lower() not in SUPPORTED_EXTS:
            continue
        try:
            text = extract_text_from_path(p)
        except Exception as e:  # missing parser, corrupt file, etc. — skip, don't abort
            log(f"  skip {p.name}: {e}")
            continue
        if not text.strip():
            continue
        chunk = f"\n\n----- {p.name} -----\n\n{text}"
        if total + len(chunk) > max_chars:
            log(f"  context size cap ({max_chars} chars) reached — skipping remaining files")
            break
        parts.append(chunk)
        total += len(chunk)
        log(f"  + {p.name} ({len(text)} chars)")
    return "".join(parts).strip()
