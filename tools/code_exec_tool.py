"""Code execution tool: runs matplotlib/mermaid code to produce diagrams."""
import io
import traceback
from pathlib import Path


def _slide_structure(slide, index: int) -> dict:
    """Extract a compact structure summary for one slide."""
    shapes_info = []
    for sh in slide.shapes:
        info = {"type": sh.shape_type.name if sh.shape_type else "UNKNOWN",
                "pos": f"{sh.left},{sh.top}", "size": f"{sh.width}x{sh.height}"}
        if sh.has_text_frame:
            text = sh.text_frame.text.strip()
            if text:
                info["text"] = text[:300]
        if hasattr(sh, "image"):
            try:
                info["image"] = f"{sh.image.content_type} ({sh.width}x{sh.height})"
            except Exception:
                info["has_image"] = True
        shapes_info.append(info)
    result = {"index": index, "shape_count": len(shapes_info), "shapes": shapes_info}
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            result["speaker_notes"] = notes[:500]
    return result


def run_pptx_code(code: str, pptx_path: str) -> dict:
    """
    Execute python-pptx code against a target .pptx file.

    The code runs with `prs` (Presentation) and `pptx_path` in scope. If the
    file doesn't exist, a blank Presentation is opened. After code finishes,
    the presentation is saved.

    Returns {"success": True, "slide_count": N, "stdout": ..., "slides_touched": [...]}
    or {"success": False, "error": "..."}.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    import os as _os
    import io as _io

    pptx_path = str(Path(pptx_path).resolve())
    Path(pptx_path).parent.mkdir(parents=True, exist_ok=True)

    if Path(pptx_path).exists():
        prs = Presentation(pptx_path)
    else:
        prs = Presentation()

    # Capture print() output from user code
    stdout_capture = io.StringIO()

    namespace = {
        "prs": prs,
        "Presentation": Presentation,
        "Inches": Inches,
        "Pt": Pt,
        "Emu": Emu,
        "RGBColor": RGBColor,
        "PP_ALIGN": PP_ALIGN,
        "MSO_SHAPE": MSO_SHAPE,
        "pptx_path": pptx_path,
        "__builtins__": {
            "print": lambda *a, **kw: stdout_capture.write(" ".join(str(x) for x in a) + "\n"),
            "range": range, "len": len, "enumerate": enumerate,
            "zip": zip, "list": list, "dict": dict, "tuple": tuple, "set": set,
            "str": str, "int": int, "float": float, "bool": bool,
            "abs": abs, "round": round, "max": max, "min": min, "sum": sum,
            "True": True, "False": False, "None": None,
            "isinstance": isinstance, "hasattr": hasattr, "getattr": getattr,
            "repr": repr, "type": type,
        },
    }

    # Snapshot slide hashes before exec to detect which slides changed
    def _slide_hash(slide):
        parts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                parts.append(sh.text_frame.text[:100])
            parts.append(f"{sh.left},{sh.top},{sh.width},{sh.height}")
        return hash(tuple(parts))

    before_hashes = {i: _slide_hash(s) for i, s in enumerate(prs.slides)}

    try:
        exec(compile(code, "<pptx-agent>", "exec"), namespace)
        prs = namespace.get("prs", prs)
        prs.save(pptx_path)

        # Detect which slides were touched
        after_hashes = {i: _slide_hash(s) for i, s in enumerate(prs.slides)}
        touched_indices = [
            i for i in range(len(prs.slides))
            if after_hashes.get(i) != before_hashes.get(i)
            or i >= len(before_hashes)
        ]

        # Return structure of touched slides so model sees result immediately
        slides_touched = [_slide_structure(prs.slides[i], i) for i in touched_indices[:3]]

        stdout_text = stdout_capture.getvalue().strip()
        result = {
            "success": True,
            "slide_count": len(prs.slides),
            "path": pptx_path,
            "slides_touched": slides_touched,
        }
        if stdout_text:
            result["stdout"] = stdout_text[:2000]
        return result
    except Exception as e:
        return {"success": False, "error": traceback.format_exc()[-1800:]}


def read_pptx_summary(pptx_path: str, slide_index: int | None = None) -> dict:
    """Open a .pptx and return a structural summary.

    If slide_index is provided, returns detailed info for just that slide.
    Otherwise returns a summary of all slides (less detail per slide).
    """
    from pptx import Presentation
    if not Path(pptx_path).exists():
        return {"error": f"No file at {pptx_path}", "slides": []}
    prs = Presentation(pptx_path)

    if slide_index is not None:
        if slide_index < 0 or slide_index >= len(prs.slides):
            return {"error": f"slide_index {slide_index} out of range (deck has {len(prs.slides)} slides)"}
        slide = prs.slides[slide_index]
        return {"slide_count": len(prs.slides), "slide": _slide_structure(slide, slide_index)}

    slides = []
    for i, slide in enumerate(prs.slides):
        slides.append(_slide_structure(slide, i))
    return {"slide_count": len(prs.slides), "slides": slides}


def run_diagram_code(code: str, output_path: str, diagram_type: str = "matplotlib") -> dict:
    """
    Execute diagram generation code in a restricted namespace.
    Returns {"success": True, "path": output_path} or {"success": False, "error": "..."}.

    Only matplotlib is supported for security. The code must save to `output_path`.
    """
    output_path = str(Path(output_path).resolve())
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)

    if diagram_type == "matplotlib":
        return _run_matplotlib(code, output_path)
    else:
        return {"success": False, "error": f"Unsupported diagram type: {diagram_type}"}


def _run_matplotlib(code: str, output_path: str) -> dict:
    """Execute matplotlib code in a sandboxed namespace."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    # Inject output_path into the namespace so code can reference it
    namespace = {
        "plt": plt,
        "np": np,
        "output_path": output_path,
        "__builtins__": {
            # Allow safe builtins only
            "print": print,
            "range": range,
            "len": len,
            "enumerate": enumerate,
            "zip": zip,
            "list": list,
            "dict": dict,
            "str": str,
            "int": int,
            "float": float,
            "bool": bool,
            "abs": abs,
            "round": round,
            "max": max,
            "min": min,
            "sum": sum,
        },
    }

    # Ensure code saves the figure
    if "savefig" not in code:
        code += f"\nplt.savefig(output_path, dpi=150, bbox_inches='tight')"
    if "plt.close" not in code:
        code += "\nplt.close('all')"

    try:
        exec(compile(code, "<diagram>", "exec"), namespace)
        if Path(output_path).exists():
            return {"success": True, "path": output_path}
        else:
            return {"success": False, "error": "Code executed but no file was saved."}
    except Exception:
        return {"success": False, "error": traceback.format_exc()[-1800:]}
